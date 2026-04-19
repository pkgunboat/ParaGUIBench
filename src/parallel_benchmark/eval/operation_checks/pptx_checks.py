"""
PPT 演示文稿（.pptx）属性检查原语。

每个函数接收文件路径和参数字典，返回标准化结果：
    {"pass": bool, "score": float 0.0~1.0, "reason": str}

依赖: python-pptx
"""

import logging
from typing import Any, Dict, Optional

from pptx import Presentation
from pptx.util import Emu

logger = logging.getLogger("eval.operation_checks.pptx")


# ------------------------------------------------------------------
# 工具函数
# ------------------------------------------------------------------

def _load_presentation(file_path: str) -> Optional[Presentation]:
    """
    安全加载 pptx 文件。

    输入:
        file_path: pptx 文件路径
    输出:
        Presentation 对象；加载失败返回 None
    """
    try:
        return Presentation(file_path)
    except Exception as exc:
        logger.error("无法打开 pptx 文件 %s: %s", file_path, exc)
        return None


def _ok(reason: str = "通过") -> Dict[str, Any]:
    return {"pass": True, "score": 1.0, "reason": reason}


def _fail(reason: str) -> Dict[str, Any]:
    return {"pass": False, "score": 0.0, "reason": reason}


def _partial(score: float, reason: str) -> Dict[str, Any]:
    return {"pass": score >= 0.5, "score": round(score, 4), "reason": reason}


# ------------------------------------------------------------------
# Transition 类型映射
# ------------------------------------------------------------------

# python-pptx 不直接暴露 transition 类型，需要解析 XML
# Transition XML tag 到人类可读名称的映射
_TRANSITION_TAG_MAP = {
    "blinds": "Blinds",
    "checker": "Checker",
    "comb": "Comb",
    "cover": "Cover",
    "cut": "Cut",
    "diamond": "Diamond",
    "dissolve": "Dissolve",
    "fade": "Fade",
    "newsflash": "Newsflash",
    "plus": "Plus",
    "pull": "Pull",
    "push": "Push",
    "random": "Random",
    "randomBar": "Random Bar",
    "split": "Split",
    "strips": "Strips",
    "uncover": "Uncover",
    "wedge": "Wedge",
    "wheel": "Wheel",
    "wipe": "Wipe",
    "zoom": "Zoom",
    # Office 2010+ 新增
    "vortex": "Vortex",
    "shred": "Shred",
    "ripple": "Ripple",
    "honeycomb": "Honeycomb",
    "glitter": "Glitter",
    "doors": "Doors",
    "window": "Window",
    "ferris": "Ferris",
    "gallery": "Gallery",
    "conveyor": "Conveyor",
    "pan": "Pan",
    "fly": "Fly",
    "curtains": "Curtains",
    "flashBulb": "Flash",
    "prism": "Prism",
    "reveal": "Reveal",
    "warp": "Warp",
}


def _get_slide_transition_type(slide) -> Optional[str]:
    """
    从幻灯片 XML 中提取 transition 类型名称。

    输入:
        slide: python-pptx Slide 对象
    输出:
        transition 类型名称字符串；无 transition 时返回 None
    """
    from lxml import etree

    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    }

    slide_xml = slide._element
    transition = slide_xml.find("p:transition", nsmap)
    if transition is None:
        return None

    # 查找 transition 子元素（即具体的效果类型）
    for child in transition:
        tag = etree.QName(child).localname
        # 在标准命名空间和 p14 命名空间中查找
        normalized = tag.lower()
        for key, name in _TRANSITION_TAG_MAP.items():
            if key.lower() == normalized:
                return name

        # 未在映射表中找到，返回原始标签名
        return tag

    return None


# ------------------------------------------------------------------
# 检查函数
# ------------------------------------------------------------------

def check_slide_transition(file_path: str, params: dict) -> dict:
    """
    检查指定幻灯片的切换（transition）效果类型。

    输入:
        file_path: pptx 文件路径
        params:
            slide_indices (list[int]): 需要检查的幻灯片索引列表（0-based）；
                为空则检查所有幻灯片
            transition_type (str): 期望的 transition 类型名称，如 "Dissolve", "Fade"
                大小写不敏感
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    slide_indices = params.get("slide_indices")
    expected_type = params.get("transition_type", "")

    if not expected_type:
        return _fail("参数缺少 transition_type")

    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")

    slides = list(prs.slides)
    if not slides:
        return _fail("演示文稿无幻灯片")

    # 确定检查范围
    if slide_indices:
        indices = [i for i in slide_indices if 0 <= i < len(slides)]
    else:
        indices = list(range(len(slides)))

    if not indices:
        return _fail(f"指定的幻灯片索引超出范围（共 {len(slides)} 张）")

    total = len(indices)
    matched = 0
    details = []

    for idx in indices:
        actual = _get_slide_transition_type(slides[idx])
        if actual and actual.lower() == expected_type.lower():
            matched += 1
        else:
            details.append(f"第{idx + 1}张: {actual or '无transition'}")

    ratio = matched / total
    if ratio >= 1.0:
        return _ok(
            f"全部 {total} 张幻灯片 transition 为 {expected_type}"
        )

    if ratio > 0:
        detail_str = "; ".join(details[:5])
        return _partial(
            ratio,
            f"transition 匹配率 {ratio:.1%}（{matched}/{total}），"
            f"不匹配: {detail_str}"
        )

    return _fail(
        f"无幻灯片 transition 为 {expected_type}（{'; '.join(details[:5])}）"
    )


def check_text_not_overflow(file_path: str, params: dict) -> dict:
    """
    检查幻灯片中的文本框内容是否在边界内（不溢出）。

    溢出判断方式：检查文本框的文字区域是否超出形状边界。
    由于 python-pptx 无法直接计算渲染后的文本高度，
    采用启发式方法——估算文本行数与文本框可用高度的比例。

    输入:
        file_path: pptx 文件路径
        params:
            slide_indices (list[int], 可选): 检查的幻灯片索引（0-based）
            line_height_pt (float): 估算的单行高度（磅），默认 18.0
            tolerance_ratio (float): 溢出容忍比例，默认 1.2（超出 20% 算溢出）
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    slide_indices = params.get("slide_indices")
    line_height_pt = params.get("line_height_pt", 18.0)
    tolerance_ratio = params.get("tolerance_ratio", 1.2)

    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")

    slides = list(prs.slides)
    if not slides:
        return _fail("演示文稿无幻灯片")

    if slide_indices:
        check_slides = [(i, slides[i]) for i in slide_indices if 0 <= i < len(slides)]
    else:
        check_slides = list(enumerate(slides))

    total_shapes = 0
    overflow_shapes = 0
    overflow_details = []

    line_height_emu = Emu(int(line_height_pt * 12700))  # 1pt = 12700 EMU

    for slide_idx, slide in check_slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            text = tf.text.strip()
            if not text:
                continue

            total_shapes += 1

            # 估算文本行数（按换行符 + 每段落字符数 / 估算每行字符数）
            shape_width_emu = shape.width
            # 粗略估算每行能放的字符数（假设平均字符宽度 ≈ 行高 * 0.6）
            avg_char_width = line_height_emu * 0.6
            chars_per_line = max(1, int(shape_width_emu / avg_char_width)) if avg_char_width > 0 else 40

            estimated_lines = 0
            for para in tf.paragraphs:
                para_text = para.text
                if not para_text:
                    estimated_lines += 1
                else:
                    estimated_lines += max(1, len(para_text) / chars_per_line)

            # 估算需要的高度
            needed_height = estimated_lines * line_height_emu
            available_height = shape.height

            if needed_height > available_height * tolerance_ratio:
                overflow_shapes += 1
                overflow_details.append(
                    f"第{slide_idx + 1}张 形状'{shape.name}': "
                    f"估算 {estimated_lines:.0f} 行，高度不足"
                )

    if total_shapes == 0:
        return _ok("无文本框需要检查")

    ratio = 1.0 - (overflow_shapes / total_shapes)
    if overflow_shapes == 0:
        return _ok(f"全部 {total_shapes} 个文本框未溢出")

    detail_str = "; ".join(overflow_details[:5])
    return _partial(
        ratio,
        f"{overflow_shapes}/{total_shapes} 个文本框疑似溢出: {detail_str}"
    )


# ------------------------------------------------------------------
# 任务专用检查函数
# ------------------------------------------------------------------

def _rects_overlap(a, b) -> bool:
    """
    判断两个矩形是否重叠。

    输入:
        a, b: 各为 (left, top, right, bottom) 四元组，单位 EMU
    输出:
        bool，True 表示存在重叠区域
    """
    a_left, a_top, a_right, a_bottom = a
    b_left, b_top, b_right, b_bottom = b
    return (a_left < b_right and a_right > b_left and
            a_top < b_bottom and a_bottom > b_top)


def check_batchppt002_bounds_overlap(file_path: str, params: dict) -> dict:
    """
    Batchoperationppt-002 专用：坐标级文本框边界与重叠检测。

    对每张幻灯片：
      1. 检查每个文本框的四个顶点是否在幻灯片边界内
      2. 检查同一张幻灯片中多个文本框之间是否存在矩形重叠

    输入:
        file_path: pptx 文件路径
        params: 忽略
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slides = list(prs.slides)

    if not slides:
        return _fail("演示文稿无幻灯片")

    total_issues = 0
    details = []

    for slide_idx, slide in enumerate(slides):
        textbox_rects = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue

            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height

            # 检查边界：是否超出幻灯片
            out_of_bounds = False
            if left < 0 or top < 0 or right > slide_width or bottom > slide_height:
                out_of_bounds = True
                total_issues += 1
                details.append(
                    f"第{slide_idx + 1}张 '{shape.name}' 超出边界"
                )

            textbox_rects.append((left, top, right, bottom, shape.name))

        # 检查重叠
        for i in range(len(textbox_rects)):
            for j in range(i + 1, len(textbox_rects)):
                r_i = textbox_rects[i][:4]
                r_j = textbox_rects[j][:4]
                if _rects_overlap(r_i, r_j):
                    total_issues += 1
                    name_i = textbox_rects[i][4]
                    name_j = textbox_rects[j][4]
                    details.append(
                        f"第{slide_idx + 1}张 '{name_i}' 与 '{name_j}' 重叠"
                    )

    if total_issues == 0:
        return _ok(
            f"全部 {len(slides)} 张幻灯片无越界/重叠"
        )

    detail_str = "; ".join(details[:8])
    # 有任意重叠或越界问题即判 0 分
    return _fail(f"{total_issues} 处问题: {detail_str}")


# ------------------------------------------------------------------
# 扩展检查函数
# ------------------------------------------------------------------

def check_ppt_has_images_or_tables(file_path: str, params: dict) -> dict:
    """
    检查PPT中是否存在图片或表格。

    输入:
        file_path: pptx 文件路径
        params:
            check_images (bool): 是否检查图片，默认 True
            check_tables (bool): 是否检查表格，默认 True
            threshold (float): 符合比例阈值，默认 0.5
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    check_images = params.get("check_images", True)
    check_tables = params.get("check_tables", True)
    threshold = params.get("threshold", 0.5)

    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")

    slides = list(prs.slides)
    if not slides:
        return _fail("演示文稿无幻灯片")

    slides_with_content = 0

    for slide in slides:
        has_image = False
        has_table = False

        if check_images:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    has_image = True
                    break

        if check_tables:
            for shape in slide.shapes:
                if shape.has_table:
                    has_table = True
                    break

        if has_image or has_table:
            slides_with_content += 1

    ratio = slides_with_content / len(slides)
    if ratio >= threshold:
        return _ok(f"{slides_with_content}/{len(slides)} 张幻灯片包含图片或表格")

    return _partial(ratio, f"仅 {slides_with_content}/{len(slides)} 张幻灯片包含图片或表格")


def check_ppt_slide_has_number(file_path: str, params: dict) -> dict:
    """
    检查PPT指定幻灯片中是否包含特定数字。

    输入:
        file_path: pptx 文件路径
        params:
            slide_index (int): 幻灯片索引（0-based），默认 0
            expected_number (int/float): 期望出现的数字
            threshold (float): 符合比例阈值，默认 0.8
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    slide_index = params.get("slide_index", 0)
    expected_number = params.get("expected_number")
    threshold = params.get("threshold", 0.8)

    if expected_number is None:
        return _fail("参数缺少 expected_number")

    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")

    slides = list(prs.slides)
    if slide_index < 0 or slide_index >= len(slides):
        return _fail(f"幻灯片索引 {slide_index} 超出范围（共 {len(slides)} 张）")

    slide = slides[slide_index]
    expected_str = str(expected_number)

    found = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if expected_str in text:
                found = True
                break

    if found:
        return _ok(f"第 {slide_index + 1} 张幻灯片包含数字 {expected_number}")

    return _fail(f"第 {slide_index + 1} 张幻灯片不包含数字 {expected_number}")
