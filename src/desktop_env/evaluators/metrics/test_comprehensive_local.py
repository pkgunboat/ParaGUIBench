"""
综合评估函数本地测试脚本。

通过 mock 缺失的第三方依赖，在本地 Mac 上直接测试新增的比较函数。

测试项:
  1. docx 原子函数（compare_docx_paragraph_styles, compare_docx_run_formatting）
  2. docx 综合评分（compare_docx_comprehensive）
  3. pptx 原子函数（extract_all_transitions, compare_pptx_transitions）
  4. pptx 综合评分（compare_pptx_comprehensive）
  5. compare_file_comprehensive 分发

注意: xlsx 综合测试依赖 table.py 的深层依赖链（utils.py → formulas, lxml.cssselect），
本地 Mac 环境缺少这些依赖，因此 xlsx 测试在服务器上运行。

用法:
    cd ubuntu_env
    python desktop_env/evaluators/metrics/test_comprehensive_local.py
"""

import importlib.util
import logging
import os
import sys
import tempfile
import types

# ============================================================
# 第一步：Mock 所有缺失的第三方依赖（必须在所有项目 import 之前）
# ============================================================

_MODULES_TO_MOCK = [
    "rapidfuzz", "rapidfuzz.fuzz",
    "easyocr",
    "odf", "odf.opendocument", "odf.text",
]

for _name in _MODULES_TO_MOCK:
    if _name not in sys.modules:
        _mod = types.ModuleType(_name)
        if _name == "rapidfuzz.fuzz":
            setattr(_mod, "ratio", lambda *a, **kw: 0.0)
        if _name == "odf.opendocument":
            setattr(_mod, "load", lambda *a, **kw: None)
        if _name == "odf.text":
            setattr(_mod, "P", type("P", (), {}))
            setattr(_mod, "Span", type("Span", (), {}))
        sys.modules[_name] = _mod

# ============================================================
# 路径设置
# ============================================================

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
UBUNTU_ENV_DIR = os.path.dirname(os.path.dirname(os.path.dirname(SCRIPT_DIR)))

if UBUNTU_ENV_DIR not in sys.path:
    sys.path.insert(0, UBUNTU_ENV_DIR)

# ============================================================
# 直接从文件导入单个模块（绕过 __init__.py 的全量导入链）
# ============================================================

def _load_module(name: str, filepath: str):
    """
    从指定文件路径加载模块，注册到 sys.modules。

    输入:
        name: 模块名
        filepath: 文件路径

    输出:
        模块对象
    """
    spec = importlib.util.spec_from_file_location(name, filepath)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


# 1) docs.py — 只需 rapidfuzz/easyocr/odf mock
docs_mod = _load_module(
    "desktop_env.evaluators.metrics.docs",
    os.path.join(SCRIPT_DIR, "docs.py"),
)

# 2) slides.py — 没有特殊外部依赖
slides_mod = _load_module(
    "desktop_env.evaluators.metrics.slides",
    os.path.join(SCRIPT_DIR, "slides.py"),
)

# 3) 为 comprehensive.py 提供 mock 的 table 模块
#    comprehensive.py 只调用 compare_table()，我们给它一个简单的 mock
_mock_table = types.ModuleType("desktop_env.evaluators.metrics.table")
setattr(_mock_table, "compare_table", lambda **kw: 1.0)  # xlsx 自比较 mock 为 1.0
sys.modules["desktop_env.evaluators.metrics.table"] = _mock_table

comprehensive_mod = _load_module(
    "desktop_env.evaluators.metrics.comprehensive",
    os.path.join(SCRIPT_DIR, "comprehensive.py"),
)

# ============================================================
# 提取需要测试的函数
# ============================================================

compare_docx_paragraph_styles = docs_mod.compare_docx_paragraph_styles
compare_docx_run_formatting = docs_mod.compare_docx_run_formatting
extract_all_transitions = slides_mod.extract_all_transitions
compare_pptx_transitions = slides_mod.compare_pptx_transitions
compare_docx_comprehensive = comprehensive_mod.compare_docx_comprehensive
compare_pptx_comprehensive = comprehensive_mod.compare_pptx_comprehensive
compare_file_comprehensive = comprehensive_mod.compare_file_comprehensive

# ============================================================
# 第三方库导入
# ============================================================

from docx import Document as DocxDocument
from docx.shared import RGBColor
from pptx import Presentation

# ============================================================
# 日志
# ============================================================

logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("test_comprehensive")

# ============================================================
# 测试辅助：创建文件
# ============================================================

def create_test_docx(path: str):
    """
    创建一个包含多种样式和格式的测试 docx 文件。

    内容: Heading 1 + 带格式的正文 + Heading 2 + 正文
    """
    doc = DocxDocument()
    doc.add_heading("测试标题一", level=1)
    p = doc.add_paragraph()
    run1 = p.add_run("普通文本，")
    run2 = p.add_run("加粗文本，")
    run2.bold = True
    run3 = p.add_run("红色文本。")
    run3.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    doc.add_heading("测试标题二", level=2)
    doc.add_paragraph("这是最后一段正文内容。")
    doc.save(path)


def create_test_pptx(path: str):
    """
    创建一个包含 2 张幻灯片的测试 pptx 文件。
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "测试演示文稿"
    slide.placeholders[1].text = "副标题内容"
    slide_layout2 = prs.slide_layouts[5]
    prs.slides.add_slide(slide_layout2)
    prs.save(path)


# ============================================================
# 测试用例
# ============================================================

def test_docx_paragraph_styles_identity():
    """Test 1: 段落样式自比较 → 1.0"""
    log.info("=" * 50)
    log.info("TEST 1: compare_docx_paragraph_styles 恒等测试")

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        create_test_docx(path)
        doc = DocxDocument(path)
        score = compare_docx_paragraph_styles(doc, doc)
        log.info("  得分: %.4f", score)
        assert score == 1.0, f"自比较应为 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_docx_paragraph_styles_diff():
    """Test 2: 修改段落样式 → < 1.0"""
    log.info("=" * 50)
    log.info("TEST 2: compare_docx_paragraph_styles 差异检测")

    with tempfile.TemporaryDirectory() as tmpdir:
        gt_path = os.path.join(tmpdir, "gt.docx")
        mod_path = os.path.join(tmpdir, "mod.docx")
        create_test_docx(gt_path)

        # 将 Heading 1 改为 Normal
        doc = DocxDocument(gt_path)
        for para in doc.paragraphs:
            if para.style.name == "Heading 1":
                para.style = doc.styles["Normal"]
                break
        doc.save(mod_path)

        doc_gt = DocxDocument(gt_path)
        doc_mod = DocxDocument(mod_path)
        score = compare_docx_paragraph_styles(doc_gt, doc_mod)
        log.info("  得分: %.4f（期望 < 1.0）", score)
        assert score < 1.0, f"改样式后应 < 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_docx_run_formatting_identity():
    """Test 3: Run 格式自比较 → 1.0"""
    log.info("=" * 50)
    log.info("TEST 3: compare_docx_run_formatting 恒等测试")

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        create_test_docx(path)
        doc = DocxDocument(path)
        score = compare_docx_run_formatting(doc, doc)
        log.info("  得分: %.4f", score)
        assert score == 1.0, f"自比较应为 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_docx_run_formatting_diff():
    """Test 4: 修改 Run 加粗 → < 1.0"""
    log.info("=" * 50)
    log.info("TEST 4: compare_docx_run_formatting 差异检测")

    with tempfile.TemporaryDirectory() as tmpdir:
        gt_path = os.path.join(tmpdir, "gt.docx")
        mod_path = os.path.join(tmpdir, "mod.docx")
        create_test_docx(gt_path)

        doc = DocxDocument(gt_path)
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.bold = not run.bold
                    break
            break
        doc.save(mod_path)

        doc_gt = DocxDocument(gt_path)
        doc_mod = DocxDocument(mod_path)
        score = compare_docx_run_formatting(doc_gt, doc_mod)
        log.info("  得分: %.4f（期望 < 1.0）", score)
        assert score < 1.0, f"改 bold 后应 < 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_docx_comprehensive_identity():
    """Test 5: docx 综合自比较 → 1.0"""
    log.info("=" * 50)
    log.info("TEST 5: compare_docx_comprehensive 恒等测试")

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        create_test_docx(path)
        score = compare_docx_comprehensive(path, path)
        log.info("  得分: %.4f", score)
        assert score == 1.0, f"自比较应为 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_docx_comprehensive_format_change():
    """Test 6: docx 仅改格式 → 综合得分 < 1.0"""
    log.info("=" * 50)
    log.info("TEST 6: compare_docx_comprehensive 格式修改检测")

    with tempfile.TemporaryDirectory() as tmpdir:
        gt_path = os.path.join(tmpdir, "gt.docx")
        mod_path = os.path.join(tmpdir, "mod.docx")
        create_test_docx(gt_path)

        doc = DocxDocument(gt_path)
        for para in doc.paragraphs:
            if para.style.name == "Heading 1":
                para.style = doc.styles["Normal"]
                break
        doc.save(mod_path)

        score = compare_docx_comprehensive(gt_path, mod_path)
        log.info("  得分: %.4f（期望 < 1.0）", score)
        assert score < 1.0, f"格式修改后综合得分应 < 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_docx_comprehensive_content_change():
    """Test 7: docx 改内容 → 综合得分 < 1.0"""
    log.info("=" * 50)
    log.info("TEST 7: compare_docx_comprehensive 内容修改检测")

    with tempfile.TemporaryDirectory() as tmpdir:
        gt_path = os.path.join(tmpdir, "gt.docx")
        mod_path = os.path.join(tmpdir, "mod.docx")
        create_test_docx(gt_path)

        doc = DocxDocument(gt_path)
        doc.add_paragraph("这是测试追加的新段落。")
        doc.save(mod_path)

        score = compare_docx_comprehensive(gt_path, mod_path)
        log.info("  得分: %.4f（期望 < 1.0）", score)
        assert score < 1.0, f"内容修改后综合得分应 < 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_pptx_transitions_identity():
    """Test 8: pptx transition 自比较 → 1.0"""
    log.info("=" * 50)
    log.info("TEST 8: compare_pptx_transitions 恒等测试")

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        create_test_pptx(path)

        trans = extract_all_transitions(path)
        log.info("  提取到的 transitions: %s", trans)

        score = compare_pptx_transitions(path, path)
        log.info("  得分: %.4f", score)
        assert score == 1.0, f"自比较应为 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_pptx_comprehensive_identity():
    """Test 9: pptx 综合自比较 → 1.0"""
    log.info("=" * 50)
    log.info("TEST 9: compare_pptx_comprehensive 恒等测试")

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        create_test_pptx(path)
        score = compare_pptx_comprehensive(path, path)
        log.info("  得分: %.4f", score)
        assert score == 1.0, f"自比较应为 1.0，实际 {score}"
    log.info("  PASS")
    return True


def test_compare_file_comprehensive_dispatch():
    """Test 10: compare_file_comprehensive 正确分发"""
    log.info("=" * 50)
    log.info("TEST 10: compare_file_comprehensive 分发测试")

    with tempfile.TemporaryDirectory() as tmpdir:
        docx_path = os.path.join(tmpdir, "test.docx")
        pptx_path = os.path.join(tmpdir, "test.pptx")

        create_test_docx(docx_path)
        create_test_pptx(pptx_path)

        for path, ext in [(docx_path, "docx"), (pptx_path, "pptx")]:
            score = compare_file_comprehensive(path, path)
            log.info("  %s 自比较: %.4f", ext, score)
            assert score == 1.0, f"{ext} 自比较应为 1.0，实际 {score}"

        # 不支持的扩展名
        txt_path = os.path.join(tmpdir, "test.txt")
        with open(txt_path, "w") as f:
            f.write("hello")
        score = compare_file_comprehensive(txt_path, txt_path)
        log.info("  txt（不支持）: %.4f", score)
        assert score == 0.0, f"不支持的扩展名应返回 0.0，实际 {score}"

    log.info("  PASS")
    return True


# ============================================================
# 主入口
# ============================================================

def main():
    """运行所有测试。"""
    tests = [
        test_docx_paragraph_styles_identity,
        test_docx_paragraph_styles_diff,
        test_docx_run_formatting_identity,
        test_docx_run_formatting_diff,
        test_docx_comprehensive_identity,
        test_docx_comprehensive_format_change,
        test_docx_comprehensive_content_change,
        test_pptx_transitions_identity,
        test_pptx_comprehensive_identity,
        test_compare_file_comprehensive_dispatch,
    ]

    passed = 0
    failed = 0
    errors = []

    for test_func in tests:
        try:
            test_func()
            passed += 1
        except AssertionError as e:
            failed += 1
            errors.append((test_func.__name__, str(e)))
            log.error("  FAIL: %s", e)
        except Exception as e:
            failed += 1
            errors.append((test_func.__name__, f"Exception: {e}"))
            log.error("  ERROR: %s", e, exc_info=True)

    log.info("")
    log.info("=" * 50)
    log.info("测试汇总: %d 通过, %d 失败 (共 %d)", passed, failed, passed + failed)
    log.info("注意: xlsx 综合测试因本地缺少依赖已跳过，需在服务器验证")
    if errors:
        for name, msg in errors:
            log.error("  FAIL: %s → %s", name, msg)
    else:
        log.info("所有测试通过!")

    return failed == 0


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
