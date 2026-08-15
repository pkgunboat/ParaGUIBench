"""OSWorld raw artifact/gold bytes 到强类型 pure metric value 的受信投影。

本模块仅在 evaluator-only 边界中处理 ``ArtifactFamilyCapture`` 和已完成
摘要、大小与媒体验证的 logical gold bytes。它不接受 host/guest 路径，
不下载 gold，不在 repr/异常中回显文件名或内容，也不将投影值
交给 Agent 或 RunStore。
"""

from __future__ import annotations

from collections.abc import Mapping
from dataclasses import dataclass
from io import BytesIO
import json
import math
from types import MappingProxyType
import warnings
import zipfile

from paraguibench.evaluation.operation.evaluator import (
    OperationEvaluationError,
    _MAX_ARCHIVE_MEMBERS,
    _MAX_ARCHIVE_UNCOMPRESSED_BYTES,
    _validate_archive_member,
)
from paraguibench.evaluation.osworld.artifact_metric_values import (
    DocumentParagraphValue,
    NormalizedRGBImageValue,
    PDFArchiveMemberValue,
    PDFArchiveValue,
    PDFTextValue,
    PresentationArtifactValue,
    PresentationParagraphValue,
    PresentationRunValue,
    PresentationShapeValue,
    PresentationSlideValue,
    SpreadsheetArtifactValue,
    SpreadsheetCellValue,
    SpreadsheetSheetValue,
)
from paraguibench.integrations.osworld.artifact_evidence_specs import (
    OSWORLD_ARTIFACT_EVIDENCE_SPECS,
    ArtifactEvidenceSpec,
    ArtifactMetricEvidenceSpec,
    ArtifactSlotEvidenceSpec,
)
from paraguibench.integrations.osworld.artifact_family_evidence import (
    LEGACY_OSWORLD_ARTIFACT_TASK_IDS,
    ArtifactFamilyCapture,
    _is_safe_member_name,
    _load_strict_json_object,
    _payload_kind_for_slot,
    _read_zip_member,
    _validate_zip_members,
    _verify_spec_digest,
)


_FIRST_SHEET_CONTRACT_ID = "sheet-data.first-sheet.v1"
_NAMED_SHEET_CONTRACT_ID = "sheet-data.named-unseen-movies.v1"
_FUZZY_SHEET_CONTRACT_ID = "sheet-fuzzy.restaurant-contacts.v1"
_SHEET_PRINT_CONTRACT_IDS = frozenset(
    {
        "grf-sheet-print.sheet1.v1",
        "supported-rate-sheet-print.sheet1.v1",
    }
)
_DOCX_CONTRACT_IDS = frozenset(
    {
        "docx-content.v1",
        "apa7-references.content-only.base-0_6.v1",
    }
)
_SPEAKER_NOTES_CONTRACT_ID = "speaker-notes.no-shape-no-bullets.v1"
_PROBLEM_PDF_CONTRACT_ID = "problem-invoice-content.v1"
_PDF_ARCHIVE_CONTRACT_ID = "pdf-chapter-archive.v1"
_SLIDE_BACKGROUND_IMAGE_CONTRACT_ID = "slide-index-1.frame-00-08.v1"
_MAX_PDF_PAGES = 2_048
_MAX_IMAGE_SOURCE_PIXELS = 10_000_000
_MAX_IMAGE_DECODED_BYTES = 64 * 1024 * 1024
_MAX_IMAGE_NORMALIZED_PIXELS = 2_500_000
_MAX_WORKBOOK_SHEETS = 64
_MAX_TABLE_CELLS = 1_000_000
_MAX_WORKBOOK_TOTAL_CELLS = _MAX_TABLE_CELLS
_MAX_TABLE_TEXT_CHARACTERS = 1_048_576
_MAX_DOCUMENT_PARAGRAPHS = 100_000
_MAX_PRESENTATION_SLIDES = 2_048
_MAX_PRESENTATION_SHAPES = 100_000
_MAX_PRESENTATION_PARAGRAPHS = 200_000
_MAX_PRESENTATION_RUNS = 500_000


class OSWorldArtifactMetricProjectionError(RuntimeError):
    """表示 raw capture/gold 无法可靠投影为 pure metric value。

    输入参数：
        code：不含路径、文件名、内容或 gold 身份的固定错误码。
    输出返回值：
        可由上层 evidence adapter 映射为 ``schema_error`` 的脱敏异常。
    """

    def __init__(self, code: str) -> None:
        """构造只携带固定错误码的投影异常。

        输入参数：
            code：稳定错误分类。
        输出返回值：
            无；初始化当前异常。
        """

        super().__init__(code)
        self.code = code


@dataclass(frozen=True, slots=True, repr=False)
class ArtifactMetricValueProjection:
    """保存一次 pure metric 调用的脱敏身份与私有强类型值。

    输入参数：
        contract_id/metric_id：固定 evidence spec 中的版本 contract 与源 metric。
        _actual_value/_gold_value：只向 evaluator-only consumer 开放的强类型值。
        _options_json：已通过严格 JSON object 校验的固定 options。
    输出返回值：
        不可变投影；repr 仅显示 contract/metric 身份。
    """

    contract_id: str
    metric_id: str
    _actual_value: object
    _gold_value: object
    _options_json: str

    def __repr__(self) -> str:
        """生成不含 artifact/gold/options 原值的安全调试文本。

        输入参数：
            无。
        输出返回值：
            仅含固定 contract 与 metric 身份的字符串。
        """

        return (
            "ArtifactMetricValueProjection("
            f"contract_id={self.contract_id!r}, metric_id={self.metric_id!r})"
        )

    def actual_value(self) -> object:
        """向 evaluator-only consumer 返回 Agent artifact 的强类型值。

        输入参数：
            无。
        输出返回值：
            已完成安全解析的不可变 actual value。
        """

        return self._actual_value

    def gold_value(self) -> object:
        """向 evaluator-only consumer 返回 verified gold 的强类型值。

        输入参数：
            无。
        输出返回值：
            已完成安全解析的不可变 gold value。
        """

        return self._gold_value

    def options(self) -> dict[str, object]:
        """返回从固定 canonical JSON 重新构造的 metric options。

        输入参数：
            无。
        输出返回值：
            新建且不与内部状态共享的 JSON object。
        """

        return _load_strict_json_object(self._options_json)


def project_verified_artifact_metric_values(
    task_id: str,
    capture: ArtifactFamilyCapture,
    *,
    verified_gold_bytes: Mapping[str, bytes],
) -> tuple[ArtifactMetricValueProjection, ...]:
    """把一个 available raw capture 与已验证 gold bytes 投影为 metric 输入。

    输入参数：
        task_id：13 个 legacy artifact family 中的 canonical task ID。
        capture：同一 task/spec 产生的 available ``ArtifactFamilyCapture``。
        verified_gold_bytes：evaluator-only resolver 已完成摘要、大小和媒体
            验证的 ``logical gold key -> bytes`` 精确闭集。
    输出返回值：
        按 slot spec metric 顺序排列的强类型投影 tuple。
    异常：
        OSWorldArtifactMetricProjectionError：task/spec/capture/gold/options/依赖或
            解析资源边界不完整；异常不回显任何原值。
    """

    spec, slot = _resolve_projection_binding(task_id, capture)
    gold_bytes = _validated_gold_bytes(
        slot,
        spec,
        verified_gold_bytes,
    )
    actual_items = capture.file_items()
    actual_member_names = capture.member_names()
    allows_empty_pdf_archive = (
        len(slot.metrics) == 1
        and slot.metrics[0].contract_id == _PDF_ARCHIVE_CONTRACT_ID
        and not actual_member_names
    )
    if (
        (not actual_items and not allows_empty_pdf_archive)
        or len(actual_items) > spec.limits.max_items
        or any(len(item) > spec.limits.max_single_item_bytes for item in actual_items)
        or sum(len(item) for item in actual_items) > spec.limits.max_total_bytes
    ):
        raise OSWorldArtifactMetricProjectionError("CAPTURE_SCHEMA_ERROR")

    projections: list[ArtifactMetricValueProjection] = []
    for metric in slot.metrics:
        try:
            options = _load_strict_json_object(metric.options_json)
        except (TypeError, ValueError):
            raise OSWorldArtifactMetricProjectionError("OPTIONS_SCHEMA_ERROR") from None
        if metric.contract_id in {
            _FIRST_SHEET_CONTRACT_ID,
            _NAMED_SHEET_CONTRACT_ID,
            _FUZZY_SHEET_CONTRACT_ID,
        }:
            actual_value, gold_value = _project_first_sheet_values(
                actual_items,
                gold_bytes,
                metric,
                spec,
            )
        elif metric.contract_id in _SHEET_PRINT_CONTRACT_IDS:
            actual_value, gold_value = _project_sheet_print_values(
                actual_items,
                gold_bytes,
                metric,
                spec,
            )
        elif metric.contract_id in _DOCX_CONTRACT_IDS:
            actual_value, gold_value = _project_docx_values(
                actual_items,
                gold_bytes,
                metric,
                spec,
            )
        elif metric.contract_id == _SPEAKER_NOTES_CONTRACT_ID:
            actual_value, gold_value = _project_pptx_values(
                actual_items,
                gold_bytes,
                metric,
                spec,
            )
        elif metric.contract_id == _PROBLEM_PDF_CONTRACT_ID:
            actual_value, gold_value = _project_pdf_values(
                actual_items,
                gold_bytes,
                metric,
                spec,
            )
        elif metric.contract_id == _PDF_ARCHIVE_CONTRACT_ID:
            actual_value, gold_value = _project_pdf_archive_values(
                actual_items,
                actual_member_names,
                gold_bytes,
                metric,
                spec,
            )
        elif metric.contract_id == _SLIDE_BACKGROUND_IMAGE_CONTRACT_ID:
            actual_value, gold_value = _project_image_values(
                actual_items,
                gold_bytes,
                metric,
                spec,
            )
        else:
            raise OSWorldArtifactMetricProjectionError("CONTRACT_NOT_PROJECTED")
        projections.append(
            ArtifactMetricValueProjection(
                contract_id=metric.contract_id,
                metric_id=metric.metric_id,
                _actual_value=actual_value,
                _gold_value=gold_value,
                _options_json=json.dumps(
                    options,
                    ensure_ascii=False,
                    allow_nan=False,
                    sort_keys=True,
                    separators=(",", ":"),
                ),
            )
        )
    return tuple(projections)


def _resolve_projection_binding(
    task_id: str,
    capture: ArtifactFamilyCapture,
) -> tuple[ArtifactEvidenceSpec, ArtifactSlotEvidenceSpec]:
    """绑定 canonical task、摘要验证 spec 与当前 raw capture 槽位。

    输入参数：
        task_id：候选 canonical task ID。
        capture：候选 raw capture。
    输出返回值：
        唯一匹配的 ``(ArtifactEvidenceSpec, ArtifactSlotEvidenceSpec)``。
    异常：
        OSWorldArtifactMetricProjectionError：任务、摘要、槽位或 capture 状态失配。
    """

    if (
        task_id not in LEGACY_OSWORLD_ARTIFACT_TASK_IDS
        or not isinstance(capture, ArtifactFamilyCapture)
        or capture.status != "available"
    ):
        raise OSWorldArtifactMetricProjectionError("CAPTURE_BINDING_ERROR")
    spec = OSWORLD_ARTIFACT_EVIDENCE_SPECS.get(task_id)
    if spec is None:
        raise OSWorldArtifactMetricProjectionError("TASK_NOT_REGISTERED")
    try:
        _verify_spec_digest(spec)
    except Exception:
        raise OSWorldArtifactMetricProjectionError("SPEC_IDENTITY_ERROR") from None
    slot = next(
        (
            candidate
            for candidate in spec.artifact_slots
            if candidate.slot_id == capture.slot_id
        ),
        None,
    )
    if (
        slot is None
        or not slot.metrics
        or capture.payload_kind != _payload_kind_for_slot(slot)
    ):
        raise OSWorldArtifactMetricProjectionError("CAPTURE_BINDING_ERROR")
    return spec, slot


def _validated_gold_bytes(
    slot: ArtifactSlotEvidenceSpec,
    spec: ArtifactEvidenceSpec,
    values: Mapping[str, bytes],
) -> Mapping[str, bytes]:
    """校验外部 gold logical key 闭集与二次字节预算。

    输入参数：
        slot/spec：已绑定的槽位和任务取证规格。
        values：候选 logical gold key 到 bytes 映射。
    输出返回值：
        key 和 bytes 均复制到新建只读映射的精确 gold 闭集。
    异常：
        OSWorldArtifactMetricProjectionError：缺 key、多 key、非 bytes 或超出 spec 预算。
    """

    expected_keys = tuple(key for metric in slot.metrics for key in metric.gold_keys)
    if (
        not isinstance(values, Mapping)
        or not expected_keys
        or len(set(expected_keys)) != len(expected_keys)
        or set(values) != set(expected_keys)
    ):
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    copied = dict(values)
    if (
        any(
            not isinstance(value, bytes)
            or len(value) > spec.limits.max_single_item_bytes
            for value in copied.values()
        )
        or sum(len(value) for value in copied.values()) > spec.limits.max_total_bytes
    ):
        raise OSWorldArtifactMetricProjectionError("GOLD_SCHEMA_ERROR")
    return MappingProxyType(copied)


def _project_first_sheet_values(
    actual_items: tuple[bytes, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[SpreadsheetArtifactValue, SpreadsheetArtifactValue]:
    """把 first-sheet contract 的单个 actual/gold XLSX 投影为 workbook 值。

    输入参数：
        actual_items：raw capture 中的文件 bytes tuple。
        gold_bytes：已二次校验的 logical gold 映射。
        metric/spec：固定 metric 身份与资源上限。
    输出返回值：
        ``(actual workbook value, gold workbook value)``。
    异常：
        OSWorldArtifactMetricProjectionError：文件数、gold key 或 XLSX 安全解析失败。
    """

    if len(actual_items) != 1 or len(metric.gold_keys) != 1:
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_item = gold_bytes.get(metric.gold_keys[0])
    if gold_item is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    return (
        _parse_xlsx_workbook(
            actual_items[0],
            spec,
            first_sheet_only=metric.contract_id == _FIRST_SHEET_CONTRACT_ID,
        ),
        _parse_xlsx_workbook(
            gold_item,
            spec,
            first_sheet_only=metric.contract_id == _FIRST_SHEET_CONTRACT_ID,
        ),
    )


def _project_sheet_print_values(
    actual_items: tuple[bytes, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[SpreadsheetArtifactValue, SpreadsheetArtifactValue]:
    """验证 XLSX 并把同 bundle 的 UTF-8 CSV sidecar 投影为 Sheet1 print 值。

    输入参数：
        actual_items：按 locator 顺序保存的 ``(xlsx, csv)`` bytes。
        gold_bytes：按 metric gold key 精确绑定的已验证 bytes。
        metric/spec：固定 metric 身份、gold 顺序与资源预算。
    输出返回值：
        只含 ``Sheet1.printed_text`` 的 actual/gold workbook 值对。
    异常：
        OSWorldArtifactMetricProjectionError：输入数量、XLSX 预检/解析、UTF-8
            或文本预算无效。
    """

    if len(actual_items) != 2 or len(metric.gold_keys) != 2:
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_xlsx = gold_bytes.get(metric.gold_keys[0])
    gold_csv = gold_bytes.get(metric.gold_keys[1])
    if gold_xlsx is None or gold_csv is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    _parse_xlsx_workbook(actual_items[0], spec)
    _parse_xlsx_workbook(gold_xlsx, spec)
    actual_text = _decode_bounded_utf8(actual_items[1], spec)
    gold_text = _decode_bounded_utf8(gold_csv, spec)
    return (
        SpreadsheetArtifactValue(
            (SpreadsheetSheetValue("Sheet1", printed_text=actual_text),)
        ),
        SpreadsheetArtifactValue(
            (SpreadsheetSheetValue("Sheet1", printed_text=gold_text),)
        ),
    )


def _decode_bounded_utf8(content: bytes, spec: ArtifactEvidenceSpec) -> str:
    """在双重字节上限内严格解码 artifact/gold UTF-8 文本。

    输入参数：
        content：待解码的原始 bytes。
        spec：提供 ``max_text_bytes`` 的固定取证规格。
    输出返回值：
        无 BOM 特别处理、严格 UTF-8 解码的字符串。
    异常：
        OSWorldArtifactMetricProjectionError：超限或存在无效 UTF-8。
    """

    if len(content) > spec.limits.max_text_bytes:
        raise OSWorldArtifactMetricProjectionError("TEXT_LIMIT_EXCEEDED")
    try:
        return content.decode("utf-8", errors="strict")
    except UnicodeDecodeError:
        raise OSWorldArtifactMetricProjectionError("TEXT_DECODE_ERROR") from None


def _project_docx_values(
    actual_items: tuple[bytes, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[DocumentParagraphValue, DocumentParagraphValue]:
    """把单个 actual/gold DOCX 安全投影为有序段落值。

    输入参数：
        actual_items：raw capture 中的单文件 bytes tuple。
        gold_bytes：已校验 logical gold 映射。
        metric/spec：固定 gold key 与 OOXML/文本资源预算。
    输出返回值：
        actual 与 gold 的 ``DocumentParagraphValue`` 对。
    异常：
        OSWorldArtifactMetricProjectionError：输入数量、gold key 或 DOCX 预检/解析无效。
    """

    if len(actual_items) != 1 or len(metric.gold_keys) != 1:
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_item = gold_bytes.get(metric.gold_keys[0])
    if gold_item is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    return (
        _parse_docx_paragraphs(actual_items[0], spec),
        _parse_docx_paragraphs(gold_item, spec),
    )


def _parse_docx_paragraphs(
    content: bytes,
    spec: ArtifactEvidenceSpec,
) -> DocumentParagraphValue:
    """预检 DOCX 并延迟使用 python-docx 提取顶层段落文本。

    输入参数：
        content：已受 raw/gold 预算约束的 DOCX bytes。
        spec：容器与文本上限的固定取证规格。
    输出返回值：
        按 ``Document.paragraphs`` 顺序保存的不可变文本投影。
    异常：
        OSWorldArtifactMetricProjectionError：OOXML 预检、依赖、段落数或文本预算失败。
    """

    _preflight_ooxml_bytes(content, spec)
    try:
        from docx import Document
    except ImportError:
        raise OSWorldArtifactMetricProjectionError("DEPENDENCY_MISSING") from None
    try:
        document = Document(BytesIO(content))
        paragraphs = tuple(paragraph.text for paragraph in document.paragraphs)
    except Exception:
        raise OSWorldArtifactMetricProjectionError("DOCX_PARSE_ERROR") from None
    if (
        len(paragraphs) > _MAX_DOCUMENT_PARAGRAPHS
        or sum(len(paragraph) for paragraph in paragraphs) > spec.limits.max_text_bytes
    ):
        raise OSWorldArtifactMetricProjectionError("DOCX_TEXT_LIMIT_EXCEEDED")
    return DocumentParagraphValue(paragraphs)


def _project_pptx_values(
    actual_items: tuple[bytes, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[PresentationArtifactValue, PresentationArtifactValue]:
    """把单个 actual/gold PPTX 投影为 speaker-notes 强类型值。

    输入参数：
        actual_items：raw capture 中的单文件 bytes tuple。
        gold_bytes：已校验 logical gold 映射。
        metric/spec：固定 gold key 与 OOXML/文本资源预算。
    输出返回值：
        actual 与 gold 的 ``PresentationArtifactValue`` 对。
    异常：
        OSWorldArtifactMetricProjectionError：输入数量、gold key 或 PPTX
            安全解析失败。
    """

    if len(actual_items) != 1 or len(metric.gold_keys) != 1:
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_item = gold_bytes.get(metric.gold_keys[0])
    if gold_item is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    return (
        _parse_pptx_presentation(actual_items[0], spec),
        _parse_pptx_presentation(gold_item, spec),
    )


def _parse_pptx_presentation(
    content: bytes,
    spec: ArtifactEvidenceSpec,
) -> PresentationArtifactValue:
    """预检 PPTX 并延迟提取源 speaker-notes contract 仍启用的字段。

    输入参数：
        content：已受 raw/gold 字节预算约束的 PPTX bytes。
        spec：提供 OOXML 容器与文本上限的固定规格。
    输出返回值：
        按 slide/shape/paragraph/run 顺序保存的不可变投影。
    异常：
        OSWorldArtifactMetricProjectionError：OOXML 预检、依赖、解析或累计
            资源预算失败。
    """

    _preflight_ooxml_bytes(content, spec)
    try:
        from pptx import Presentation
    except ImportError:
        raise OSWorldArtifactMetricProjectionError("DEPENDENCY_MISSING") from None
    try:
        presentation = Presentation(BytesIO(content))
        if len(presentation.slides) > _MAX_PRESENTATION_SLIDES:
            raise OSWorldArtifactMetricProjectionError("PPTX_STRUCTURE_LIMIT_EXCEEDED")
        slides: list[PresentationSlideValue] = []
        shape_count = paragraph_count = run_count = text_characters = 0
        for slide in presentation.slides:
            shape_count += len(slide.shapes)
            if shape_count > _MAX_PRESENTATION_SHAPES:
                raise OSWorldArtifactMetricProjectionError(
                    "PPTX_STRUCTURE_LIMIT_EXCEEDED"
                )
            notes_text = slide.notes_slide.notes_text_frame.text
            if not isinstance(notes_text, str):
                raise OSWorldArtifactMetricProjectionError("PPTX_SCHEMA_ERROR")
            text_characters += len(notes_text)
            shapes: list[PresentationShapeValue] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    shapes.append(PresentationShapeValue(None))
                    continue
                shape_text = shape.text
                if not isinstance(shape_text, str) or not hasattr(shape, "text_frame"):
                    raise OSWorldArtifactMetricProjectionError("PPTX_SCHEMA_ERROR")
                text_characters += len(shape_text)
                paragraphs: list[PresentationParagraphValue] = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_count += 1
                    if paragraph_count > _MAX_PRESENTATION_PARAGRAPHS:
                        raise OSWorldArtifactMetricProjectionError(
                            "PPTX_STRUCTURE_LIMIT_EXCEEDED"
                        )
                    paragraph_text = paragraph.text
                    if not isinstance(paragraph_text, str):
                        raise OSWorldArtifactMetricProjectionError("PPTX_SCHEMA_ERROR")
                    text_characters += len(paragraph_text)
                    runs: list[PresentationRunValue] = []
                    for run in paragraph.runs:
                        run_count += 1
                        if run_count > _MAX_PRESENTATION_RUNS:
                            raise OSWorldArtifactMetricProjectionError(
                                "PPTX_STRUCTURE_LIMIT_EXCEEDED"
                            )
                        runs.append(_project_presentation_run(run))
                    paragraphs.append(
                        PresentationParagraphValue(
                            text=paragraph_text,
                            alignment=_project_office_enum(paragraph.alignment),
                            level=int(paragraph.level),
                            runs=tuple(runs),
                        )
                    )
                shapes.append(
                    PresentationShapeValue(
                        text=shape_text,
                        paragraphs=tuple(paragraphs),
                    )
                )
            if text_characters > spec.limits.max_text_bytes:
                raise OSWorldArtifactMetricProjectionError("PPTX_TEXT_LIMIT_EXCEEDED")
            # 源 evaluator 调用 ``fill.background()``；python-pptx 该 setter
            # 固定返回 None，因而源 contract 中此字段始终为 None。
            slides.append(
                PresentationSlideValue(
                    background_color=None,
                    notes_text=notes_text,
                    shapes=tuple(shapes),
                )
            )
        return PresentationArtifactValue(tuple(slides))
    except OSWorldArtifactMetricProjectionError:
        raise
    except Exception:
        raise OSWorldArtifactMetricProjectionError("PPTX_PARSE_ERROR") from None


def _project_presentation_run(run: object) -> PresentationRunValue:
    """把 python-pptx run 收紧为仅含允许标量的强类型值。

    输入参数：
        run：python-pptx 解析后的文本 run 对象。
    输出返回值：
        字体名、字号、粗斜体、RGB、下划线和删除线标量。
    异常：
        OSWorldArtifactMetricProjectionError：解析器返回了不可固定的属性。
    """

    try:
        font = run.font
        font_name = font.name
        font_size = None if font.size is None else int(font.size)
        color_rgb = None
        try:
            rgb = font.color.rgb
        except (AttributeError, TypeError, ValueError):
            rgb = None
        if rgb is not None:
            color_rgb = str(rgb)
        strike = font._element.attrib.get("strike", "noStrike")
        projected = PresentationRunValue(
            font_name=font_name,
            font_size=font_size,
            bold=font.bold,
            italic=font.italic,
            color_rgb=color_rgb,
            underline=_project_office_enum(font.underline),
            strike=strike,
        )
    except Exception:
        raise OSWorldArtifactMetricProjectionError("PPTX_SCHEMA_ERROR") from None
    if (
        (font_name is not None and not isinstance(font_name, str))
        or (projected.bold is not None and type(projected.bold) is not bool)
        or (projected.italic is not None and type(projected.italic) is not bool)
        or not isinstance(strike, str)
    ):
        raise OSWorldArtifactMetricProjectionError("PPTX_SCHEMA_ERROR")
    return projected


def _project_office_enum(value: object) -> str | int | bool | None:
    """把 Office enum 或原生标量投影为跨 parser 实例可比较的值。

    输入参数：
        value：python-pptx 返回的对齐或下划线值。
    输出返回值：
        ``None``、精确 bool/int/str，或 enum 的稳定整数/名称。
    异常：
        OSWorldArtifactMetricProjectionError：该值无法可靠投影。
    """

    if value is None or type(value) in {str, int, bool}:
        return value
    if isinstance(value, int):
        return int(value)
    name = getattr(value, "name", None)
    if isinstance(name, str):
        return name
    raise OSWorldArtifactMetricProjectionError("PPTX_SCHEMA_ERROR")


def _project_pdf_values(
    actual_items: tuple[bytes, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[PDFTextValue, PDFTextValue]:
    """把单个实际 PDF 与已验证 gold PDF 投影为文本值。

    输入参数：
        actual_items：raw capture 中的 PDF bytes tuple。
        gold_bytes：已验证的 logical gold key 到 PDF bytes 映射。
        metric/spec：固定 metric 身份和资源预算。
    输出返回值：
        ``(actual PDFTextValue, gold PDFTextValue)``。
    异常：
        OSWorldArtifactMetricProjectionError：文件数或 gold 绑定不完整。
    """

    if len(actual_items) != 1 or len(metric.gold_keys) != 1:
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_item = gold_bytes.get(metric.gold_keys[0])
    if gold_item is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    return (
        _parse_pdf_text(actual_items[0], spec),
        _parse_pdf_text(gold_item, spec),
    )


def _parse_pdf_text(
    content: bytes,
    spec: ArtifactEvidenceSpec,
) -> PDFTextValue:
    """使用延迟导入的受控 PDF parser 提取有界页面文本。

    输入参数：
        content：已受 raw/gold 总量预算限制的 PDF bytes。
        spec：提供最大文本字节数的固定取证规格。
    输出返回值：
        按页面顺序拼接并去除首尾空白的 ``PDFTextValue``。
    异常：
        OSWorldArtifactMetricProjectionError：依赖缺失、签名、加密、页数、
            文本大小或解析结果无效。
    """

    if not content.startswith(b"%PDF-") or b"%%EOF" not in content[-65_536:]:
        raise OSWorldArtifactMetricProjectionError("PDF_PARSE_ERROR")
    try:
        from pypdf import PdfReader
    except ImportError:
        raise OSWorldArtifactMetricProjectionError("DEPENDENCY_MISSING") from None
    stream = BytesIO(content)
    try:
        reader = PdfReader(stream, strict=True)
        if reader.is_encrypted:
            raise OSWorldArtifactMetricProjectionError("PDF_ENCRYPTED_REJECTED")
        pages = reader.pages
        if len(pages) > _MAX_PDF_PAGES:
            raise OSWorldArtifactMetricProjectionError("PDF_PAGE_LIMIT_EXCEEDED")
        extracted: list[str] = []
        text_characters = 0
        for page in pages:
            text = page.extract_text()
            if text is None:
                text = ""
            if not isinstance(text, str):
                raise OSWorldArtifactMetricProjectionError("PDF_PARSE_ERROR")
            text_characters += len(text)
            if text_characters > spec.limits.max_text_bytes:
                raise OSWorldArtifactMetricProjectionError("PDF_TEXT_LIMIT_EXCEEDED")
            extracted.append(text)
        return PDFTextValue("".join(extracted).strip())
    except OSWorldArtifactMetricProjectionError:
        raise
    except Exception:
        raise OSWorldArtifactMetricProjectionError("PDF_PARSE_ERROR") from None
    finally:
        stream.close()


def _project_pdf_archive_values(
    actual_items: tuple[bytes, ...],
    actual_names: tuple[str, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[PDFArchiveValue, PDFArchiveValue]:
    """投影实际 PDF bundle 与已验证 gold ZIP 的成员闭集。

    输入参数：
        actual_items/actual_names：同序实际 PDF bytes 与顶层成员名。
        gold_bytes：已验证的 logical gold key 到 ZIP bytes 映射。
        metric/spec：固定 metric 身份和容器资源预算。
    输出返回值：
        ``(actual PDFArchiveValue, gold PDFArchiveValue)``。
    异常：
        OSWorldArtifactMetricProjectionError：数量、名称、gold 或归档无效。
    """

    if len(metric.gold_keys) != 1 or len(actual_items) != len(actual_names):
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_item = gold_bytes.get(metric.gold_keys[0])
    if gold_item is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    return (
        _build_pdf_archive_value(actual_names, actual_items, spec),
        _parse_gold_pdf_archive(gold_item, spec),
    )


def _build_pdf_archive_value(
    names: tuple[str, ...],
    items: tuple[bytes, ...],
    spec: ArtifactEvidenceSpec,
) -> PDFArchiveValue:
    """把顶层 PDF 名称与内容闭集构造为强类型归档值。

    输入参数：
        names/items：同序顶层成员名与 PDF bytes tuple。
        spec：提供最大成员数量和 PDF 文本预算的固定规格。
    输出返回值：
        按名称排序的 ``PDFArchiveValue``。
    异常：
        OSWorldArtifactMetricProjectionError：类型、数量、名称或 PDF 无效。
    """

    if (
        not isinstance(names, tuple)
        or not isinstance(items, tuple)
        or len(names) != len(items)
        or len(names) > spec.limits.max_items
        or len(set(names)) != len(names)
        or any(
            not _is_safe_member_name(name) or "/" in name or not name.endswith(".pdf")
            for name in names
        )
    ):
        raise OSWorldArtifactMetricProjectionError("PDF_ARCHIVE_SCHEMA_ERROR")
    members = tuple(
        PDFArchiveMemberValue(name, _parse_pdf_text(payload, spec))
        for name, payload in sorted(zip(names, items), key=lambda item: item[0])
    )
    return PDFArchiveValue(members)


def _parse_gold_pdf_archive(
    content: bytes,
    spec: ArtifactEvidenceSpec,
) -> PDFArchiveValue:
    """安全解包已验证 gold ZIP 并投影其顶层 PDF 闭集。

    输入参数：
        content：已完成外层摘要与大小验证的 ZIP bytes。
        spec：提供成员数、单成员及总展开量预算的固定规格。
    输出返回值：
        按名称排序的 ``PDFArchiveValue``。
    异常：
        OSWorldArtifactMetricProjectionError：ZIP 结构、成员或 PDF 无效。
    """

    try:
        with zipfile.ZipFile(BytesIO(content), mode="r") as archive:
            members = _validate_zip_members(
                archive,
                max_entries=min(
                    spec.limits.max_container_entries,
                    _MAX_ARCHIVE_MEMBERS,
                ),
                max_single_member_bytes=spec.limits.max_single_item_bytes,
                max_expanded_bytes=min(
                    spec.limits.max_container_expanded_bytes,
                    _MAX_ARCHIVE_UNCOMPRESSED_BYTES,
                ),
            )
            names: list[str] = []
            items: list[bytes] = []
            for member in sorted(
                members.values(),
                key=lambda candidate: candidate.filename,
            ):
                _validate_archive_member(member)
                name = member.filename
                if (
                    member.is_dir()
                    or not _is_safe_member_name(name)
                    or "/" in name
                    or not name.endswith(".pdf")
                ):
                    raise ValueError("gold PDF archive member schema 无效")
                names.append(name)
                items.append(
                    _read_zip_member(
                        archive,
                        member,
                        max_bytes=spec.limits.max_single_item_bytes,
                    )
                )
    except (OperationEvaluationError, ValueError, zipfile.BadZipFile):
        raise OSWorldArtifactMetricProjectionError("PDF_ARCHIVE_PARSE_ERROR") from None
    except Exception:
        raise OSWorldArtifactMetricProjectionError("PDF_ARCHIVE_PARSE_ERROR") from None
    return _build_pdf_archive_value(tuple(names), tuple(items), spec)


def _project_image_values(
    actual_items: tuple[bytes, ...],
    gold_bytes: Mapping[str, bytes],
    metric: ArtifactMetricEvidenceSpec,
    spec: ArtifactEvidenceSpec,
) -> tuple[NormalizedRGBImageValue, NormalizedRGBImageValue]:
    """联合归一化实际图与已验证 gold 图并投影 RGB/HSV 像素。

    输入参数：
        actual_items：raw capture 中的编码图像 bytes tuple。
        gold_bytes：已验证的 logical gold key 到编码图像 bytes 映射。
        metric/spec：固定 metric 身份与已绑定的取证资源预算。
    输出返回值：
        共同尺寸下的 ``(actual value, gold value)``。
    异常：
        OSWorldArtifactMetricProjectionError：数量、依赖、编码、尺寸、模式或
            解码结果不完整。
    """

    if (
        not isinstance(spec, ArtifactEvidenceSpec)
        or len(actual_items) != 1
        or len(metric.gold_keys) != 1
    ):
        raise OSWorldArtifactMetricProjectionError("METRIC_INPUT_ARITY_ERROR")
    gold_item = gold_bytes.get(metric.gold_keys[0])
    if gold_item is None:
        raise OSWorldArtifactMetricProjectionError("GOLD_BINDING_ERROR")
    try:
        from PIL import Image, ImageMode
    except ImportError:
        raise OSWorldArtifactMetricProjectionError("DEPENDENCY_MISSING") from None

    actual_rgb = None
    gold_rgb = None
    actual_normalized = None
    gold_normalized = None
    actual_hsv = None
    gold_hsv = None
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            actual_rgb = _load_bounded_rgb_image(
                actual_items[0],
                Image,
                ImageMode,
            )
            gold_rgb = _load_bounded_rgb_image(
                gold_item,
                Image,
                ImageMode,
            )
        common_width = min(actual_rgb.width, gold_rgb.width)
        common_height = min(actual_rgb.height, gold_rgb.height)
        common_pixels = common_width * common_height
        if (
            common_width <= 0
            or common_height <= 0
            or common_pixels > _MAX_IMAGE_NORMALIZED_PIXELS
        ):
            raise OSWorldArtifactMetricProjectionError("IMAGE_SIZE_LIMIT_EXCEEDED")
        common_size = (common_width, common_height)
        actual_normalized = actual_rgb.resize(
            common_size,
            resample=Image.Resampling.LANCZOS,
        )
        gold_normalized = gold_rgb.resize(
            common_size,
            resample=Image.Resampling.LANCZOS,
        )
        actual_hsv = actual_normalized.convert("HSV")
        gold_hsv = gold_normalized.convert("HSV")
        actual_rgb_bytes = actual_normalized.tobytes()
        gold_rgb_bytes = gold_normalized.tobytes()
        actual_hsv_bytes = actual_hsv.tobytes()
        gold_hsv_bytes = gold_hsv.tobytes()
        expected_length = common_pixels * 3
        if any(
            len(payload) != expected_length
            for payload in (
                actual_rgb_bytes,
                gold_rgb_bytes,
                actual_hsv_bytes,
                gold_hsv_bytes,
            )
        ):
            raise OSWorldArtifactMetricProjectionError("IMAGE_PARSE_ERROR")
        return (
            NormalizedRGBImageValue(
                width=common_width,
                height=common_height,
                rgb_pixels=actual_rgb_bytes,
                hsv_pixels=actual_hsv_bytes,
            ),
            NormalizedRGBImageValue(
                width=common_width,
                height=common_height,
                rgb_pixels=gold_rgb_bytes,
                hsv_pixels=gold_hsv_bytes,
            ),
        )
    except OSWorldArtifactMetricProjectionError:
        raise
    except Exception:
        raise OSWorldArtifactMetricProjectionError("IMAGE_PARSE_ERROR") from None
    finally:
        for image in (
            actual_hsv,
            gold_hsv,
            actual_normalized,
            gold_normalized,
            actual_rgb,
            gold_rgb,
        ):
            if image is not None:
                image.close()


def _load_bounded_rgb_image(
    content: bytes,
    image_module: object,
    image_mode_module: object,
) -> object:
    """在像素和解码内存预算内把编码图像转换为 RGB。

    输入参数：
        content：已受外层 artifact 大小预算限制的编码图像 bytes。
        image_module：延迟导入的 Pillow ``Image`` 模块。
        image_mode_module：延迟导入的 Pillow ``ImageMode`` 模块。
    输出返回值：
        调用方负责关闭的 Pillow RGB image。
    异常：
        OSWorldArtifactMetricProjectionError：头部尺寸、像素数、模式或估算
            解码内存超限。
    """

    stream = BytesIO(content)
    raw_image = None
    try:
        raw_image = image_module.open(stream)
        width = raw_image.width
        height = raw_image.height
        if (
            type(width) is not int
            or type(height) is not int
            or width <= 0
            or height <= 0
            or width * height > _MAX_IMAGE_SOURCE_PIXELS
        ):
            raise OSWorldArtifactMetricProjectionError("IMAGE_SIZE_LIMIT_EXCEEDED")
        descriptor = image_mode_module.getmode(raw_image.mode)
        bands = descriptor.bands
        type_string = descriptor.typestr
        if (
            not isinstance(bands, tuple)
            or not bands
            or not isinstance(type_string, str)
            or not type_string
        ):
            raise OSWorldArtifactMetricProjectionError("IMAGE_MODE_REJECTED")
        try:
            sample_bytes = int(type_string[-1])
        except (TypeError, ValueError):
            raise OSWorldArtifactMetricProjectionError("IMAGE_MODE_REJECTED") from None
        if sample_bytes not in {1, 2, 4, 8}:
            raise OSWorldArtifactMetricProjectionError("IMAGE_MODE_REJECTED")
        if width * height * len(bands) * sample_bytes > _MAX_IMAGE_DECODED_BYTES:
            raise OSWorldArtifactMetricProjectionError("IMAGE_SIZE_LIMIT_EXCEEDED")
        raw_image.load()
        return raw_image.convert("RGB")
    finally:
        if raw_image is not None:
            raw_image.close()
        stream.close()


def _parse_xlsx_workbook(
    content: bytes,
    spec: ArtifactEvidenceSpec,
    *,
    first_sheet_only: bool = False,
) -> SpreadsheetArtifactValue:
    """预检并延迟使用 openpyxl 把 XLSX bytes 投影为表格值。

    输入参数：
        content：已受 raw/gold 预算限制的 XLSX bytes。
        spec：提供容器成员数、单项、展开与 XML 上限的固定取证规格。
        first_sheet_only：真时只投影 workbook 首页，不解析协议无关的后续页。
    输出返回值：
        按 workbook 顺序保存、将首行投影为 columns 的不可变强类型值。
    异常：
        OSWorldArtifactMetricProjectionError：OOXML 边界、依赖、工作簿形状或单元格无效。
    """

    _preflight_ooxml_bytes(content, spec)
    try:
        import openpyxl
    except ImportError:
        raise OSWorldArtifactMetricProjectionError("DEPENDENCY_MISSING") from None
    stream = BytesIO(content)
    try:
        workbook = openpyxl.load_workbook(
            stream,
            read_only=True,
            data_only=True,
            keep_links=False,
        )
    except Exception:
        raise OSWorldArtifactMetricProjectionError("XLSX_PARSE_ERROR") from None
    try:
        if len(workbook.worksheets) > _MAX_WORKBOOK_SHEETS:
            raise OSWorldArtifactMetricProjectionError("XLSX_SHEET_LIMIT_EXCEEDED")
        sheets: list[SpreadsheetSheetValue] = []
        text_characters = 0
        worksheets = (
            workbook.worksheets[:1] if first_sheet_only else workbook.worksheets
        )
        total_cells = 0
        for worksheet in worksheets:
            if (
                type(worksheet.max_row) is not int
                or type(worksheet.max_column) is not int
                or worksheet.max_row < 1
                or worksheet.max_column < 1
                or worksheet.max_row * worksheet.max_column > _MAX_TABLE_CELLS
            ):
                raise OSWorldArtifactMetricProjectionError("XLSX_SHAPE_LIMIT_EXCEEDED")
            total_cells += worksheet.max_row * worksheet.max_column
            if total_cells > _MAX_WORKBOOK_TOTAL_CELLS:
                raise OSWorldArtifactMetricProjectionError(
                    "XLSX_WORKBOOK_CELL_LIMIT_EXCEEDED"
                )
        for worksheet in worksheets:
            raw_rows: list[tuple[object, ...]] = []
            projected_cells: list[SpreadsheetCellValue] = []
            for raw_row in worksheet.iter_rows(values_only=True):
                row = tuple(_normalize_spreadsheet_scalar(value) for value in raw_row)
                text_characters += sum(
                    len(value) for value in row if isinstance(value, str)
                )
                if text_characters > _MAX_TABLE_TEXT_CHARACTERS:
                    raise OSWorldArtifactMetricProjectionError(
                        "XLSX_TEXT_LIMIT_EXCEEDED"
                    )
                raw_rows.append(row)
            for cell_row in worksheet.iter_rows():
                for cell in cell_row:
                    normalized = _normalize_spreadsheet_scalar(cell.value)
                    projected_cells.append(
                        SpreadsheetCellValue(
                            coordinate=cell.coordinate,
                            value=str(normalized),
                        )
                    )
            while raw_rows and all(value is None for value in raw_rows[-1]):
                raw_rows.pop()
            max_column = max(
                (
                    index + 1
                    for row in raw_rows
                    for index, value in enumerate(row)
                    if value is not None
                ),
                default=0,
            )
            normalized_rows = tuple(row[:max_column] for row in raw_rows)
            columns = normalized_rows[0] if normalized_rows else ()
            data_rows = normalized_rows[1:] if normalized_rows else ()
            sheets.append(
                SpreadsheetSheetValue(
                    name=worksheet.title,
                    columns=columns,
                    rows=data_rows,
                    cells=tuple(projected_cells),
                )
            )
        if not sheets:
            raise OSWorldArtifactMetricProjectionError("XLSX_SCHEMA_ERROR")
        return SpreadsheetArtifactValue(tuple(sheets))
    except OSWorldArtifactMetricProjectionError:
        raise
    except Exception:
        raise OSWorldArtifactMetricProjectionError("XLSX_PARSE_ERROR") from None
    finally:
        workbook.close()
        stream.close()


def _normalize_spreadsheet_scalar(value: object) -> object:
    """把 openpyxl 单元格值收紧为 pure table contract 允许的标量。

    输入参数：
        value：openpyxl ``values_only`` 返回的候选单元格值。
    输出返回值：
        ``None`` 或精确 bool/int/finite-float/str。
    异常：
        OSWorldArtifactMetricProjectionError：日期、错误对象、非有限数或其他未固定类型。
    """

    if value is None or type(value) in {bool, int, str}:
        return value
    if type(value) is float and math.isfinite(value):
        return value
    raise OSWorldArtifactMetricProjectionError("XLSX_CELL_SCHEMA_ERROR")


def _preflight_ooxml_bytes(content: bytes, spec: ArtifactEvidenceSpec) -> None:
    """在任何 Office parser 导入前复用现有 ZIP/member 门禁检查 OOXML bytes。

    输入参数：
        content：待预检的受限 OOXML bytes。
        spec：提供该任务容器资源上限的固定取证规格。
    输出返回值：
        无；所有成员路径、类型、加密、压缩比、宏、总量与主动 XML
        门禁通过时正常返回。
    异常：
        OSWorldArtifactMetricProjectionError：容器无效或任一安全/资源边界失败。
    """

    try:
        with zipfile.ZipFile(BytesIO(content), mode="r") as archive:
            members = _validate_zip_members(
                archive,
                max_entries=min(
                    spec.limits.max_container_entries,
                    _MAX_ARCHIVE_MEMBERS,
                ),
                max_single_member_bytes=spec.limits.max_single_item_bytes,
                max_expanded_bytes=min(
                    spec.limits.max_container_expanded_bytes,
                    _MAX_ARCHIVE_UNCOMPRESSED_BYTES,
                ),
            )
            casefolded_names: set[str] = set()
            for member in members.values():
                _validate_archive_member(member)
                folded = member.filename.casefold()
                if folded in casefolded_names:
                    raise OperationEvaluationError("ARCHIVE_DUPLICATE_MEMBER_REJECTED")
                casefolded_names.add(folded)
                if member.filename.lower().endswith((".xml", ".rels")):
                    payload = _read_zip_member(
                        archive,
                        member,
                        max_bytes=min(
                            spec.limits.max_text_bytes,
                            spec.limits.max_single_item_bytes,
                        ),
                    )
                    lowered = payload.replace(b"\x00", b"").lower()
                    if b"<!doctype" in lowered or b"<!entity" in lowered:
                        raise OperationEvaluationError("ARCHIVE_ACTIVE_XML_REJECTED")
    except (OperationEvaluationError, ValueError, zipfile.BadZipFile):
        raise OSWorldArtifactMetricProjectionError("OOXML_PREFLIGHT_ERROR") from None
    except Exception:
        raise OSWorldArtifactMetricProjectionError("OOXML_PREFLIGHT_ERROR") from None


__all__ = [
    "ArtifactMetricValueProjection",
    "OSWorldArtifactMetricProjectionError",
    "project_verified_artifact_metric_values",
]
