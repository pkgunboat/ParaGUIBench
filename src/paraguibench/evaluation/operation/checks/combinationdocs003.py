"""CombinationDocs-003 源表格相对图像检查。"""

from __future__ import annotations

from collections.abc import Sequence
from functools import partial
import io
from pathlib import Path
from typing import Any


_SOURCE_WORKBOOK = "McDonalds_Monthly_Data.xlsx"
_OUTPUT_PRESENTATION = "McDonalds_powerpoint_report.pptx"
_SOURCE_SHEET = "Monthly Data"
_SOURCE_RANGE = "A1:F16"
_TARGET_SLIDE_INDEX = 2
_SLIDE_SIZE = (12_192_000, 6_858_000)
_EXPECTED_SLIDE_TEXT = (
    ("Mcdonald’s annual report", "2026"),
    (
        "Annual total",
        "Revenue (USD) :1398815$\nProfit (USD):316340$\nCustomers:35140",
    ),
    ("Here is chart of McDonald's - Regional Monthly Performance 2026",),
    (
        "Jan data",
        "Revenue (USD) :136690$\nProfit (USD):47109$\nCustomers:3602",
    ),
    (
        "March data",
        ("Revenue (USD) :81835$\nProfit (USD):-5823$\nAvg. Transaction (USD):43.16"),
    ),
)
_BASELINE_PLACEHOLDERS = (
    (
        ("CENTER_TITLE", 0, 1_524_000, 1_122_363, 9_144_000, 2_387_600),
        ("SUBTITLE", 1, 1_524_000, 3_602_038, 9_144_000, 1_655_762),
    ),
    (
        ("TITLE", 0, 838_200, 365_125, 10_515_600, 1_325_563),
        ("OBJECT", 1, 838_200, 1_825_625, 10_515_600, 4_351_338),
    ),
) + (
    (
        ("TITLE", 0, 838_200, 365_125, 10_515_600, 1_325_563),
        ("OBJECT", 1, 838_200, 1_825_625, 10_515_600, 4_351_338),
    ),
) * 3
_CONTENT_BOUNDS = (838_200, 1_825_625, 10_515_600, 4_351_338)
_MIN_INSERT_AXIS_NUMERATOR = 7
_MIN_INSERT_AXIS_DENOMINATOR = 10
_GRID_EXTENT_TOLERANCE_EMU = 6
_COLUMN_WIDTHS = (190, 210, 185, 190, 150, 255)
_ROW_HEIGHTS = (48, 20, 48) + (36,) * 12 + (42,)
_NORMALIZED_SIZE = (sum(_COLUMN_WIDTHS), sum(_ROW_HEIGHTS))
_MAX_IMAGE_PIXELS = 20_000_000
_MIN_PICTURE_WIDTH = 640
_MIN_PICTURE_HEIGHT = 240
_PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_RELATIONSHIP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_IMAGE_RELATIONSHIP_TYPE = f"{_RELATIONSHIP_NS}/image"
_PNG_CONTENT_TYPE = "image/png"


def _result(passed: bool, reason: str) -> dict[str, Any]:
    """构造不包含文件路径、表格值或图像摘要的内部结果。

    输入参数：
        passed：任务约束是否全部满足；reason：固定内部原因码。
    输出返回值：
        兼容 Operation 聚合器的 ``pass/score/reason`` 映射。
    """

    return {"pass": passed, "score": 1.0 if passed else 0.0, "reason": reason}


def _normalized_text(value: str) -> str:
    """规范化 PPT 文本中的换行与行尾空白。

    输入参数：
        value：python-pptx 返回的单个形状文本。
    输出返回值：
        保留可见内容和换行顺序、去除行尾空白的字符串。
    """

    return "\n".join(line.rstrip() for line in value.strip().splitlines())


def _read_source_rows(workbook_path: Path) -> tuple[tuple[object, ...], ...] | None:
    """读取固定可见范围，并拒绝隐藏、公式或结构漂移。

    输入参数：
        workbook_path：任务根内的正式源 XLSX 路径。
    输出返回值：
        ``A1:F16`` 的不可变值元组；无法可靠解析时返回 ``None``。
    """

    try:
        from openpyxl import load_workbook

        workbook = load_workbook(
            workbook_path,
            data_only=False,
            read_only=False,
            keep_links=False,
        )
    except Exception:
        return None
    try:
        if workbook.sheetnames != [_SOURCE_SHEET]:
            return None
        worksheet = workbook[_SOURCE_SHEET]
        if (
            worksheet.sheet_state != "visible"
            or worksheet.calculate_dimension() != _SOURCE_RANGE
            or tuple(str(item) for item in worksheet.merged_cells.ranges) != ("A1:F1",)
        ):
            return None
        for index in range(1, 17):
            if worksheet.row_dimensions[index].hidden:
                return None
        for column in "ABCDEF":
            if worksheet.column_dimensions[column].hidden:
                return None
        rows = tuple(
            tuple(cell.value for cell in row)
            for row in worksheet.iter_rows(
                min_row=1,
                max_row=16,
                min_col=1,
                max_col=6,
            )
        )
        if any(
            isinstance(value, str) and value.startswith("=")
            for row in rows
            for value in row
        ):
            return None
        return rows
    except Exception:
        return None
    finally:
        workbook.close()


def _render_source_projection(rows: Sequence[Sequence[object]]) -> Any:
    """将源单元格值渲染为固定尺寸的 canonical 基准投影。

    输入参数：
        rows：已通过结构安全检查的 16x6 源表格值。
    输出返回值：
        Pillow RGB Image；只在 evaluator 内存中使用，不落盘。
    """

    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", _NORMALIZED_SIZE, "white")
    draw = ImageDraw.Draw(image)
    normal = ImageFont.load_default(size=19)
    bold = ImageFont.load_default(size=20)
    y = 0
    for row_index, (row, height) in enumerate(zip(rows, _ROW_HEIGHTS, strict=True)):
        x = 0
        for column_index, (value, width) in enumerate(
            zip(row, _COLUMN_WIDTHS, strict=True)
        ):
            if row_index == 2:
                fill = "#4472C4"
                foreground = "white"
            elif row_index == 15:
                fill = "#E7E6E6"
                foreground = "black"
            else:
                fill = "white"
                foreground = "black"
            if row_index not in {0, 1}:
                draw.rectangle(
                    (x, y, x + width, y + height),
                    fill=fill,
                    outline="black",
                )
            if value is not None:
                text = str(value)
                font = bold if row_index in {0, 2, 15} else normal
                box = draw.textbbox((0, 0), text, font=font)
                text_width = box[2] - box[0]
                text_height = box[3] - box[1]
                if row_index == 0:
                    text_x = (_NORMALIZED_SIZE[0] - text_width) // 2
                elif column_index == 0 or row_index == 2:
                    text_x = x + (width - text_width) // 2
                else:
                    text_x = x + width - text_width - 8
                draw.text(
                    (text_x, y + (height - text_height) // 2 - box[1]),
                    text,
                    fill=foreground,
                    font=font,
                )
            x += width
        y += height
    return image


def _load_candidate_image(blob: bytes) -> Any | None:
    """完整解码内嵌图片，并以固定像素上限防止解压放大。

    输入参数：
        blob：PPTX 内嵌 picture relationship 的原始字节。
    输出返回值：
        全不透明 RGB PNG Pillow Image；格式、尺寸或解码无效时
        返回 ``None``。
    """

    try:
        from PIL import Image

        if not blob.startswith(b"\x89PNG\r\n\x1a\n"):
            return None
        with Image.open(io.BytesIO(blob)) as image:
            width, height = image.size
            if (
                image.format != "PNG"
                or image.mode != "RGB"
                or width < _MIN_PICTURE_WIDTH
                or height < _MIN_PICTURE_HEIGHT
                or width * height > _MAX_IMAGE_PIXELS
            ):
                return None
            image.load()
            if "transparency" in image.info:
                return None
            return image.copy()
    except Exception:
        return None


def _visual_projection_matches(expected: Any, actual: Any) -> bool:
    """精确比较源投影与候选图片的尺寸和 RGB 像素身份。

    输入参数：
        expected：从固定源 XLSX 即时生成的投影；actual：幻灯片内嵌图。
    输出返回值：
        尺寸和 RGB 字节逐字节完全一致时返回 ``True``。
    """

    return bool(expected.size == actual.size and expected.tobytes() == actual.tobytes())


def _source_cell_text(value: object) -> str:
    """将源单元格值转换为 PowerPoint 表格的可见文本。

    输入参数：
        value：openpyxl 返回的标量单元格值。
    输出返回值：
        空值对应空字符串，其他值按可见文本精确转换。
    """

    return "" if value is None else str(value)


def _rgb_hex(color_format: Any) -> str | None:
    """安全读取 python-pptx 颜色对象的显式 RGB 值。

    输入参数：
        color_format：单元格填充或 run font 的 ColorFormat。
    输出返回值：
        六位大写 RGB；主题色、无色或读取失败时返回 ``None``。
    """

    try:
        rgb = color_format.rgb
    except (AttributeError, ValueError):
        return None
    return None if rgb is None else str(rgb).upper()


def _proportional_dimensions(total: int, weights: Sequence[int]) -> tuple[int, ...]:
    """把固定源投影比例换算为总 EMU 内的整数尺寸。

    输入参数：
        total：table shape 水平或垂直总 EMU；weights：源投影列宽或行高权重。
    输出返回值：
        按比例向下取整、最后一项吸收余数且总和精确为 ``total``
        的尺寸元组。
    """

    weight_total = sum(weights)
    values = [total * weight // weight_total for weight in weights[:-1]]
    values.append(total - sum(values))
    return tuple(values)


def _native_table_matches_source(
    table: Any,
    rows: Sequence[Sequence[object]],
    *,
    shape_width: int,
    shape_height: int,
) -> bool:
    """逐格校验 native PPT table 的值、合并关系与关键样式。

    输入参数：
        table：python-pptx Table；rows：固定源 XLSX 16x6 可见值；
        shape_width/shape_height：graphicFrame 外层几何尺寸。
    输出返回值：
        值矩阵精确一致、A1:F1 唯一合并，且标题/表头/合计样式
        与源表格关键语义一致时返回 ``True``。
    """

    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.enum.text import PP_ALIGN

    if len(table.rows) != 16 or len(table.columns) != 6:
        return False
    actual_widths = tuple(column.width for column in table.columns)
    actual_heights = tuple(row.height for row in table.rows)
    if (
        abs(sum(actual_widths) - shape_width) > _GRID_EXTENT_TOLERANCE_EMU
        or abs(sum(actual_heights) - shape_height) > _GRID_EXTENT_TOLERANCE_EMU
    ):
        return False
    expected_widths = _proportional_dimensions(sum(actual_widths), _COLUMN_WIDTHS)
    expected_heights = _proportional_dimensions(sum(actual_heights), _ROW_HEIGHTS)
    if any(
        abs(actual - expected) > _GRID_EXTENT_TOLERANCE_EMU
        for actual, expected in zip(actual_widths, expected_widths, strict=True)
    ) or any(
        abs(actual - expected) > _GRID_EXTENT_TOLERANCE_EMU
        for actual, expected in zip(actual_heights, expected_heights, strict=True)
    ):
        return False
    if not table.cell(0, 0).is_merge_origin:
        return False
    if any(not table.cell(0, column).is_spanned for column in range(1, 6)):
        return False
    for row_index in range(1, 16):
        if any(
            table.cell(row_index, column).is_merge_origin
            or table.cell(row_index, column).is_spanned
            for column in range(6)
        ):
            return False
    for row_index, source_row in enumerate(rows):
        for column_index, source_value in enumerate(source_row):
            expected = _source_cell_text(source_value)
            if _normalized_text(table.cell(row_index, column_index).text) != expected:
                return False
            if expected:
                cell = table.cell(row_index, column_index)
                horizontal_margin = (cell.margin_left or 0) + (cell.margin_right or 0)
                vertical_margin = (cell.margin_top or 0) + (cell.margin_bottom or 0)
                if (
                    horizontal_margin * 2 > actual_widths[column_index]
                    or vertical_margin * 2 > actual_heights[row_index]
                ):
                    return False
    title_paragraph = table.cell(0, 0).text_frame.paragraphs[0]
    if (
        title_paragraph.alignment != PP_ALIGN.CENTER
        or not title_paragraph.runs
        or not all(run.font.bold is True for run in title_paragraph.runs)
        or any(
            run.font.size is None or run.font.size.pt < 8.0
            for run in title_paragraph.runs
        )
        or any(_rgb_hex(run.font.color) != "000000" for run in title_paragraph.runs)
    ):
        return False
    for column_index in range(6):
        header = table.cell(2, column_index)
        paragraph = header.text_frame.paragraphs[0]
        if (
            header.fill.type != MSO_FILL_TYPE.SOLID
            or _rgb_hex(header.fill.fore_color) != "4472C4"
            or paragraph.alignment != PP_ALIGN.CENTER
            or not paragraph.runs
            or not all(run.font.bold is True for run in paragraph.runs)
            or not all(_rgb_hex(run.font.color) == "FFFFFF" for run in paragraph.runs)
            or any(
                run.font.size is None or run.font.size.pt < 8.0
                for run in paragraph.runs
            )
        ):
            return False
        total = table.cell(15, column_index)
        if (
            total.fill.type != MSO_FILL_TYPE.SOLID
            or _rgb_hex(total.fill.fore_color) != "E7E6E6"
            or not total.text_frame.paragraphs[0].runs
            or not all(
                run.font.bold is True for run in total.text_frame.paragraphs[0].runs
            )
            or any(
                run.font.size is None or run.font.size.pt < 8.0
                for run in total.text_frame.paragraphs[0].runs
            )
            or any(
                _rgb_hex(run.font.color) != "000000"
                for run in total.text_frame.paragraphs[0].runs
            )
        ):
            return False
    for row_index in range(3, 15):
        for column_index in range(6):
            cell = table.cell(row_index, column_index)
            paragraph = cell.text_frame.paragraphs[0]
            expected_alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.RIGHT
            if paragraph.alignment != expected_alignment or not paragraph.runs:
                return False
            if cell.fill.type == MSO_FILL_TYPE.SOLID and _rgb_hex(
                cell.fill.fore_color
            ) not in {None, "FFFFFF"}:
                return False
            for run in paragraph.runs:
                if run.font.bold is True:
                    return False
                if run.font.size is None or run.font.size.pt < 8.0:
                    return False
                if _rgb_hex(run.font.color) != "000000":
                    return False
    return True


def _shape_is_hidden_or_transformed(shape: Any) -> bool:
    """检查形状是否通过隐藏、翻转或透明变换逃逸可见性约束。

    输入参数：
        shape：python-pptx 形状对象。
    输出返回值：
        存在 hidden 属性、非零旋转、水平/垂直翻转或 alpha 变换时
        返回 ``True``。
    """

    if shape.rotation not in {None, 0, 0.0}:
        return True
    for element in shape._element.iter():
        raw_tag = str(element.tag)
        local_name = raw_tag.rsplit("}", 1)[-1].rsplit(":", 1)[-1]
        if local_name == "cNvPr" and element.get("hidden", "0").lower() in {
            "1",
            "true",
        }:
            return True
        if local_name in {"alpha", "alphaMod", "alphaModFix", "alphaOff", "alphaRepl"}:
            return True
        if local_name == "xfrm" and (
            element.get("flipH", "0").lower() in {"1", "true"}
            or element.get("flipV", "0").lower() in {"1", "true"}
        ):
            return True
    return False


def _qname(namespace: str, local_name: str) -> str:
    """构造 OOXML 精确 Clark QName。

    输入参数：
        namespace：固定 OOXML namespace URI；local_name：局部名。
    输出返回值：
        ``{namespace}local`` 格式的精确标记。
    """

    return f"{{{namespace}}}{local_name}"


def _plain_picture_xml(shape: Any) -> bool:
    """按明确允许列表校验无效果、无裁剪的标准 ``p:pic``。

    输入参数：
        shape：第 3 页的唯一 picture 形状。
    输出返回值：
        XML 层次、子元素和属性只包含受审计的嵌入图、拉伸填充、
        坐标与矩形几何时返回 ``True``；任何 effect/ext/外链都失败。
    """

    root = shape._element
    p = partial(_qname, _PRESENTATION_NS)
    a = partial(_qname, _DRAWING_NS)
    r = partial(_qname, _RELATIONSHIP_NS)
    if root.tag != p("pic") or root.attrib:
        return False
    root_children = tuple(root)
    if tuple(item.tag for item in root_children) != (
        p("nvPicPr"),
        p("blipFill"),
        p("spPr"),
    ):
        return False
    non_visual, blip_fill, shape_properties = root_children
    if non_visual.attrib or tuple(item.tag for item in non_visual) != (
        p("cNvPr"),
        p("cNvPicPr"),
        p("nvPr"),
    ):
        return False
    common, picture_properties, application_properties = tuple(non_visual)
    if (
        set(common.attrib) - {"id", "name", "descr", "title"}
        or not {"id", "name"}.issubset(common.attrib)
        or tuple(common)
        or application_properties.attrib
        or tuple(application_properties)
        or picture_properties.attrib
    ):
        return False
    picture_property_children = tuple(picture_properties)
    if len(picture_property_children) > 1:
        return False
    if picture_property_children:
        locks = picture_property_children[0]
        if (
            locks.tag != a("picLocks")
            or set(locks.attrib) != {"noChangeAspect"}
            or locks.attrib["noChangeAspect"].lower() not in {"1", "true"}
            or tuple(locks)
        ):
            return False
    if blip_fill.attrib or tuple(item.tag for item in blip_fill) != (
        a("blip"),
        a("stretch"),
    ):
        return False
    blip, stretch = tuple(blip_fill)
    if set(blip.attrib) != {r("embed")} or tuple(blip) or stretch.attrib:
        return False
    stretch_children = tuple(stretch)
    if (
        len(stretch_children) != 1
        or stretch_children[0].tag != a("fillRect")
        or stretch_children[0].attrib
        or tuple(stretch_children[0])
    ):
        return False
    if shape_properties.attrib or tuple(item.tag for item in shape_properties) != (
        a("xfrm"),
        a("prstGeom"),
    ):
        return False
    transform, geometry = tuple(shape_properties)
    if transform.attrib or tuple(item.tag for item in transform) != (
        a("off"),
        a("ext"),
    ):
        return False
    offset, extent = tuple(transform)
    if (
        set(offset.attrib) != {"x", "y"}
        or tuple(offset)
        or set(extent.attrib) != {"cx", "cy"}
        or tuple(extent)
    ):
        return False
    geometry_children = tuple(geometry)
    return bool(
        geometry.attrib == {"prst": "rect"}
        and len(geometry_children) == 1
        and geometry_children[0].tag == a("avLst")
        and not geometry_children[0].attrib
        and not tuple(geometry_children[0])
    )


def _internal_png_blob(shape: Any, slide: Any) -> bytes | None:
    """绑定 ``r:embed`` 到 slide3 唯一的内部 PNG image relationship。

    输入参数：
        shape：已通过 raw ``p:pic`` 允许列表的图片形状；
        slide：包含该形状的第 3 页 slide。
    输出返回值：
        关系类型、TargetMode、part 路径、ContentType 和 blob 全部一致时
        返回内嵌 PNG 原始字节；任一失配返回 ``None``。
    """

    try:
        blip = shape._element.blipFill.blip
        relationship_id = blip.get(_qname(_RELATIONSHIP_NS, "embed"))
        if not relationship_id:
            return None
        relationship = slide.part.rels[relationship_id]
        if relationship.reltype != _IMAGE_RELATIONSHIP_TYPE or relationship.is_external:
            return None
        image_relationships = tuple(
            item
            for item in slide.part.rels.values()
            if item.reltype == _IMAGE_RELATIONSHIP_TYPE
        )
        if len(image_relationships) != 1 or image_relationships[0] is not relationship:
            return None
        target_part = relationship.target_part
        if (
            target_part.content_type != _PNG_CONTENT_TYPE
            or str(target_part.partname).lower().rsplit(".", 1)[-1] != "png"
        ):
            return None
        blob = target_part.blob
        if not isinstance(blob, bytes) or shape.image.blob != blob:
            return None
        return blob
    except Exception:
        return None


def _placeholder_matches(shape: Any, expected: tuple[object, ...]) -> bool:
    """校验原始幻灯片占位符的类型、索引与几何闭包。

    输入参数：
        shape：待校验形状；expected：固定 placeholder 类型名、索引与
            ``left/top/width/height``。
    输出返回值：
        所有原始结构属性精确一致时返回 ``True``。
    """

    type_name, index, left, top, width, height = expected
    return bool(
        shape.is_placeholder
        and shape.placeholder_format.type.name == type_name
        and shape.placeholder_format.idx == index
        and (shape.left, shape.top, shape.width, shape.height)
        == (left, top, width, height)
        and not _shape_is_hidden_or_transformed(shape)
    )


def _insert_shape_is_visible(shape: Any) -> bool:
    """验证新插入表格/图片位于第 3 页内容区且每轴足够。

    输入参数：
        shape：第 3 页 z-order 最后的唯一插入形状。
    输出返回值：
        完全位于原内容占位区、宽高每轴均不低于 70%，且无
        隐藏或变形通道时返回 ``True``。
    """

    left, top, width, height = _CONTENT_BOUNDS
    right = left + width
    bottom = top + height
    shape_right = shape.left + shape.width
    shape_bottom = shape.top + shape.height
    return bool(
        shape.left >= left
        and shape.top >= top
        and shape_right <= right
        and shape_bottom <= bottom
        and shape.width * _MIN_INSERT_AXIS_DENOMINATOR
        >= width * _MIN_INSERT_AXIS_NUMERATOR
        and shape.height * _MIN_INSERT_AXIS_DENOMINATOR
        >= height * _MIN_INSERT_AXIS_NUMERATOR
        and not _shape_is_hidden_or_transformed(shape)
    )


def _presentation_inserted_content(
    presentation_path: Path,
    rows: Sequence[Sequence[object]],
) -> tuple[bytes, int, int] | bool | None:
    """校验演示文稿基线与插入页，返回受审计的插入类型。

    输入参数：
        presentation_path：Agent 修改后的固定同名 PPTX 路径；
        rows：从固定源 XLSX 读取的 16x6 可见值。
    输出返回值：
        第 3 页唯一内嵌图的 ``(字节, shape 宽, shape 高)``，
        或在 native table 逐格通过时返回 ``True``；其他结构漂移
        返回 ``None``。
    """

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(presentation_path)
    except Exception:
        return None
    if (
        len(presentation.slides) != 5
        or (presentation.slide_width, presentation.slide_height) != _SLIDE_SIZE
    ):
        return None
    target_blobs: list[tuple[bytes, int, int]] = []
    matching_table_count = 0
    for slide_index, slide in enumerate(presentation.slides):
        if slide._element.get("show", "1").lower() in {"0", "false"}:
            return None
        expected_placeholders = _BASELINE_PLACEHOLDERS[slide_index]
        expected_shape_count = 3 if slide_index == _TARGET_SLIDE_INDEX else 2
        if len(slide.shapes) != expected_shape_count:
            return None
        if any(
            not _placeholder_matches(slide.shapes[index], expected)
            for index, expected in enumerate(expected_placeholders)
        ):
            return None
        visible_text = tuple(
            _normalized_text(shape.text)
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and _normalized_text(shape.text)
        )
        if visible_text != _EXPECTED_SLIDE_TEXT[slide_index]:
            return None
        if slide_index != _TARGET_SLIDE_INDEX:
            continue
        inserted_shape = slide.shapes[-1]
        if not _insert_shape_is_visible(inserted_shape):
            return None
        if inserted_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if not _plain_picture_xml(inserted_shape):
                return None
            if any(
                value not in {None, 0, 0.0}
                for value in (
                    inserted_shape.crop_left,
                    inserted_shape.crop_top,
                    inserted_shape.crop_right,
                    inserted_shape.crop_bottom,
                )
            ):
                return None
            picture_blob = _internal_png_blob(inserted_shape, slide)
            if picture_blob is None:
                return None
            target_blobs.append(
                (picture_blob, inserted_shape.width, inserted_shape.height)
            )
        elif getattr(inserted_shape, "has_table", False):
            if not _native_table_matches_source(
                inserted_shape.table,
                rows,
                shape_width=inserted_shape.width,
                shape_height=inserted_shape.height,
            ):
                return None
            matching_table_count += 1
        else:
            return None
    if matching_table_count == 1 and not target_blobs:
        return True
    if matching_table_count or len(target_blobs) != 1:
        return None
    return target_blobs[0]


def check_combinationdocs003_source_table_insert(
    root_path: str,
    params: dict[str, object],
) -> dict[str, Any]:
    """相对固定源 XLSX 评价 003 第 3 页表格内容。

    输入参数：
        root_path：已由 Operation 预检冻结的单 Attempt artifact 根；
        params：canonical 规则参数，003 v1 必须是空映射。
    输出返回值：
        仅在源范围可靠、PPT 关键内容未破坏，且第 3 页唯一
        native table 或 canonical PNG 通过完整语义合同时返回满分。
    """

    if params:
        return _result(False, "configuration_invalid")
    root = Path(root_path)
    if not root.is_dir() or root.is_symlink():
        return _result(False, "artifact_root_invalid")
    workbook_path = root / _SOURCE_WORKBOOK
    presentation_path = root / _OUTPUT_PRESENTATION
    if any(
        not path.is_file() or path.is_symlink()
        for path in (workbook_path, presentation_path)
    ):
        return _result(False, "required_artifact_invalid")
    rows = _read_source_rows(workbook_path)
    if rows is None:
        return _result(False, "source_projection_invalid")
    inserted_content = _presentation_inserted_content(presentation_path, rows)
    if inserted_content is None:
        return _result(False, "presentation_contract_mismatch")
    if inserted_content is True:
        return _result(True, "passed")
    try:
        expected = _render_source_projection(rows)
        picture_blob, shape_width, shape_height = inserted_content
        actual = _load_candidate_image(picture_blob)
        if actual is None:
            return _result(False, "picture_invalid")
        if shape_width * actual.height != shape_height * actual.width:
            return _result(False, "picture_geometry_mismatch")
        canonical_png = io.BytesIO()
        expected.save(canonical_png, format="PNG", optimize=True)
        if picture_blob != canonical_png.getvalue():
            return _result(False, "source_picture_bytes_mismatch")
        projection_matches = _visual_projection_matches(expected, actual)
    except (ImportError, OSError, ValueError):
        return _result(False, "projection_unavailable")
    if not projection_matches:
        return _result(False, "source_picture_mismatch")
    return _result(True, "passed")


COMBINATIONDOCS003_CHECKS = {
    "check_combinationdocs003_source_table_insert": (
        check_combinationdocs003_source_table_insert
    ),
}


__all__ = [
    "COMBINATIONDOCS003_CHECKS",
    "check_combinationdocs003_source_table_insert",
]
