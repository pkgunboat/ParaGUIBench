"""从最终旧源机械提取的 Operation 检查最小传递闭包。"""

from __future__ import annotations
import logging
from typing import Any, Dict, Optional
from pptx import Presentation

logger = logging.getLogger("eval.operation_checks.pptx")


def _load_presentation(file_path: str) -> Optional[Presentation]:
    """
    安全加载 pptx 文件。

    输入:
        file_path: pptx 文件路径
    输出:
        Presentation 对象；加载失败返回 None
    """
    try:
        return Presentation(file_path)
    except Exception:
        logger.error("Operation 检查内部事件")
        return None


def _ok(reason: str = "通过") -> Dict[str, Any]:
    """功能：执行传递闭包内的 _ok 检查步骤。

    输入参数：
        由函数签名定义；调用值仅在 evaluator 内存中使用。
    输出返回值：
        由函数签名定义；公开聚合层会删除路径、内容、gold 与原始 reason。"""
    return {"pass": True, "score": 1.0, "reason": reason}


def _fail(reason: str) -> Dict[str, Any]:
    """功能：执行传递闭包内的 _fail 检查步骤。

    输入参数：
        由函数签名定义；调用值仅在 evaluator 内存中使用。
    输出返回值：
        由函数签名定义；公开聚合层会删除路径、内容、gold 与原始 reason。"""
    return {"pass": False, "score": 0.0, "reason": reason}


def _partial(score: float, reason: str) -> Dict[str, Any]:
    """功能：执行传递闭包内的 _partial 检查步骤。

    输入参数：
        由函数签名定义；调用值仅在 evaluator 内存中使用。
    输出返回值：
        由函数签名定义；公开聚合层会删除路径、内容、gold 与原始 reason。"""
    return {"pass": score >= 1.0 - 1e-09, "score": round(score, 4), "reason": reason}


def _config_error(reason: str) -> Dict[str, Any]:
    """功能：执行传递闭包内的 _config_error 检查步骤。

    输入参数：
        由函数签名定义；调用值仅在 evaluator 内存中使用。
    输出返回值：
        由函数签名定义；公开聚合层会删除路径、内容、gold 与原始 reason。

    旧源语义说明：
    评价器配置错误（缺参数等）：score=-1 哨兵，由上层冒泡为 evaluator_error。"""
    return {"pass": False, "score": -1.0, "status": "evaluator_error", "reason": reason}


_TRANSITION_TAG_MAP = {
    "blinds": "Blinds",
    "checker": "Checker",
    "comb": "Comb",
    "cover": "Cover",
    "cut": "Cut",
    "diamond": "Diamond",
    "dissolve": "Dissolve",
    "fade": "Fade",
    "newsflash": "Newsflash",
    "plus": "Plus",
    "pull": "Pull",
    "push": "Push",
    "random": "Random",
    "randomBar": "Random Bar",
    "split": "Split",
    "strips": "Strips",
    "uncover": "Uncover",
    "wedge": "Wedge",
    "wheel": "Wheel",
    "wipe": "Wipe",
    "zoom": "Zoom",
    "vortex": "Vortex",
    "shred": "Shred",
    "ripple": "Ripple",
    "honeycomb": "Honeycomb",
    "glitter": "Glitter",
    "doors": "Doors",
    "window": "Window",
    "ferris": "Ferris",
    "gallery": "Gallery",
    "conveyor": "Conveyor",
    "pan": "Pan",
    "fly": "Fly",
    "curtains": "Curtains",
    "flashBulb": "Flash",
    "prism": "Prism",
    "reveal": "Reveal",
    "warp": "Warp",
}


def _get_slide_transition_type(slide) -> Optional[str]:
    """
    从幻灯片 XML 中提取 transition 类型名称。

    输入:
        slide: python-pptx Slide 对象
    输出:
        transition 类型名称字符串；无 transition 时返回 None
    """
    from lxml import etree

    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    }
    slide_xml = slide._element
    transition = slide_xml.find("p:transition", nsmap)
    if transition is None:
        return None
    for child in transition:
        tag = etree.QName(child).localname
        normalized = tag.lower()
        for key, name in _TRANSITION_TAG_MAP.items():
            if key.lower() == normalized:
                return name
        return tag
    return None


def check_slide_transition(file_path: str, params: dict) -> dict:
    """
    检查指定幻灯片的切换（transition）效果类型。

    输入:
        file_path: pptx 文件路径
        params:
            slide_indices (list[int]): 需要检查的幻灯片索引列表（0-based）；
                为空则检查所有幻灯片
            transition_type (str): 期望的 transition 类型名称，如 "Dissolve", "Fade"
                大小写不敏感
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    slide_indices = params.get("slide_indices")
    expected_type = params.get("transition_type", "")
    if not expected_type:
        return _config_error("参数缺少 transition_type")
    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")
    slides = list(prs.slides)
    if not slides:
        return _fail("演示文稿无幻灯片")
    if slide_indices:
        indices = [i for i in slide_indices if 0 <= i < len(slides)]
    else:
        indices = list(range(len(slides)))
    if not indices:
        return _fail(f"指定的幻灯片索引超出范围（共 {len(slides)} 张）")
    total = len(indices)
    matched = 0
    details = []
    for idx in indices:
        actual = _get_slide_transition_type(slides[idx])
        if actual and actual.lower() == expected_type.lower():
            matched += 1
        else:
            details.append(f"第{idx + 1}张: {actual or '无transition'}")
    ratio = matched / total
    if ratio >= 1.0:
        return _ok(f"全部 {total} 张幻灯片 transition 为 {expected_type}")
    if ratio > 0:
        detail_str = "; ".join(details[:5])
        return _partial(
            ratio,
            f"transition 匹配率 {ratio:.1%}（{matched}/{total}），不匹配: {detail_str}",
        )
    return _fail(f"无幻灯片 transition 为 {expected_type}（{'; '.join(details[:5])}）")


def _rects_overlap(a, b) -> bool:
    """
    判断两个矩形是否重叠。

    输入:
        a, b: 各为 (left, top, right, bottom) 四元组，单位 EMU
    输出:
        bool，True 表示存在重叠区域
    """
    a_left, a_top, a_right, a_bottom = a
    b_left, b_top, b_right, b_bottom = b
    return (
        a_left < b_right
        and a_right > b_left
        and (a_top < b_bottom)
        and (a_bottom > b_top)
    )


def check_batchppt002_bounds_overlap(file_path: str, params: dict) -> dict:
    """
    Batchoperationppt-002 专用：坐标级文本框边界与重叠检测。

    对每张幻灯片：
      1. 检查每个文本框的四个顶点是否在幻灯片边界内
      2. 检查同一张幻灯片中多个文本框之间是否存在矩形重叠

    输入:
        file_path: pptx 文件路径
        params: 忽略
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slides = list(prs.slides)
    if not slides:
        return _fail("演示文稿无幻灯片")
    total_issues = 0
    details = []
    for slide_idx, slide in enumerate(slides):
        textbox_rects = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue
            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height
            if left < 0 or top < 0 or right > slide_width or (bottom > slide_height):
                total_issues += 1
                details.append(f"第{slide_idx + 1}张 '{shape.name}' 超出边界")
            textbox_rects.append((left, top, right, bottom, shape.name))
        for i in range(len(textbox_rects)):
            for j in range(i + 1, len(textbox_rects)):
                r_i = textbox_rects[i][:4]
                r_j = textbox_rects[j][:4]
                if _rects_overlap(r_i, r_j):
                    total_issues += 1
                    name_i = textbox_rects[i][4]
                    name_j = textbox_rects[j][4]
                    details.append(f"第{slide_idx + 1}张 '{name_i}' 与 '{name_j}' 重叠")
    if total_issues == 0:
        return _ok(f"全部 {len(slides)} 张幻灯片无越界/重叠")
    detail_str = "; ".join(details[:8])
    return _fail(f"{total_issues} 处问题: {detail_str}")


def check_ppt_has_images_or_tables(file_path: str, params: dict) -> dict:
    """
    检查PPT中是否存在图片或表格。

    输入:
        file_path: pptx 文件路径
        params:
            check_images (bool): 是否检查图片，默认 True
            check_tables (bool): 是否检查表格，默认 True
            threshold (float): 符合比例阈值，默认 0.5
    输出:
        {"pass": bool, "score": float, "reason": str}
    """
    check_images = params.get("check_images", True)
    check_tables = params.get("check_tables", True)
    threshold = params.get("threshold", 0.5)
    prs = _load_presentation(file_path)
    if prs is None:
        return _fail(f"无法打开文件: {file_path}")
    slides = list(prs.slides)
    if not slides:
        return _fail("演示文稿无幻灯片")
    slides_with_content = 0
    for slide in slides:
        has_image = False
        has_table = False
        if check_images:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    has_image = True
                    break
        if check_tables:
            for shape in slide.shapes:
                if shape.has_table:
                    has_table = True
                    break
        if has_image or has_table:
            slides_with_content += 1
    ratio = slides_with_content / len(slides)
    if ratio >= threshold:
        return _ok(f"{slides_with_content}/{len(slides)} 张幻灯片包含图片或表格")
    return _partial(
        ratio, f"仅 {slides_with_content}/{len(slides)} 张幻灯片包含图片或表格"
    )


PPTX_CHECKS = {
    "check_slide_transition": check_slide_transition,
    "check_batchppt002_bounds_overlap": check_batchppt002_bounds_overlap,
    "check_ppt_has_images_or_tables": check_ppt_has_images_or_tables,
}
