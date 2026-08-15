"""CombinationDocs-003 源表格相对评价的行为测试。"""

from __future__ import annotations

import hashlib
import json
from pathlib import Path

import pytest

import paraguibench.evaluation.operation.evaluator as operation_evaluator_module
from paraguibench.evaluation.operation import OPERATION_TASK_RULES
from paraguibench.evaluation.operation.checks.combinationdocs003 import (
    check_combinationdocs003_source_table_insert,
)


_SOURCE_ROWS = (
    ("McDonald's - Regional Monthly Performance 2026", None, None, None, None, None),
    (None, None, None, None, None, None),
    (
        "Month",
        "Revenue (USD)",
        "Cost (USD)",
        "Profit (USD)",
        "Customers",
        "Avg. Transaction (USD)",
    ),
    ("January", 136690, 89581, 47109, 1895, 72.13),
    ("February", 88892, 97911, -9019, 3602, 24.68),
    ("March", 81835, 87658, -5823, 1896, 43.16),
    ("April", 137192, 106950, 30242, 1621, 84.63),
    ("May", 85061, 97268, -12207, 3602, 23.61),
    ("June", 110155, 73006, 37149, 3650, 30.18),
    ("July", 138887, 86245, 52642, 3069, 45.25),
    ("August", 100348, 62479, 37869, 4167, 24.08),
    ("September", 112723, 77772, 34951, 2341, 48.15),
    ("October", 135645, 109001, 26644, 3896, 34.82),
    ("November", 139448, 110760, 28688, 2885, 48.34),
    ("December", 131939, 83844, 48095, 2516, 52.44),
    ("TOTAL", 1398815, 1082475, 316340, 35140, 39.81),
)
_REPO_ROOT = Path(__file__).resolve().parents[2]
_TASK_ID = "Operation-FileOperate-CombinationDocs-003"
_COLUMN_WIDTHS = (190, 210, 185, 190, 150, 255)
_ROW_HEIGHTS = (48, 20, 48) + (36,) * 12 + (42,)


def _allocate_fixture_dimensions(
    total: int, weights: tuple[int, ...]
) -> tuple[int, ...]:
    """按测试源投影比例分配 PowerPoint 表格尺寸。

    输入参数：
        total：表格总 EMU；weights：源投影的列宽或行高权重。
    输出返回值：
        按比例向下取整且总和精确等于 ``total`` 的尺寸元组。
    """

    weight_total = sum(weights)
    dimensions = [total * weight // weight_total for weight in weights[:-1]]
    dimensions.append(total - sum(dimensions))
    return tuple(dimensions)


def _load_task() -> dict[str, object]:
    """读取 003 canonical task 供公开 Operation evaluator 测试使用。

    输入参数：
        无。
    输出返回值：
        从仓库 ``benchmark/tasks`` 读取的 JSON 映射。
    """

    return json.loads(
        (
            _REPO_ROOT
            / "benchmark/tasks/Operation-FileOperate-CombinationDocs-003.json"
        ).read_text(encoding="utf-8")
    )


def _write_source_workbook(path: Path) -> None:
    """写入与 003 正式源表格语义一致的合成工作簿。

    输入参数：
        path：待写入的 ``McDonalds_Monthly_Data.xlsx`` 路径。
    输出返回值：
        无；写入 ``Monthly Data!A1:F16`` 的值、合并、填充与边框。
    """

    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Monthly Data"
    for row in _SOURCE_ROWS:
        worksheet.append(row)
    worksheet.merge_cells("A1:F1")
    worksheet["A1"].font = Font(name="Arial", size=14, bold=True)
    worksheet["A1"].alignment = Alignment(horizontal="center")
    thin = Side(style="thin", color="000000")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for cell in worksheet[3]:
        cell.fill = PatternFill("solid", fgColor="4472C4")
        cell.font = Font(name="Arial", size=11, bold=True, color="FFFFFF")
        cell.alignment = Alignment(horizontal="center")
        cell.border = border
    for row in worksheet.iter_rows(min_row=4, max_row=16, min_col=1, max_col=6):
        for index, cell in enumerate(row, start=1):
            cell.border = border
            cell.alignment = Alignment(horizontal="left" if index == 1 else "right")
    for cell in worksheet[16]:
        cell.fill = PatternFill("solid", fgColor="E7E6E6")
        cell.font = Font(name="Arial", size=11, bold=True)
    worksheet.column_dimensions["A"].width = 18
    workbook.save(path)
    workbook.close()


def _fixture_table_picture(
    path: Path,
    rows: tuple[tuple[object, ...], ...] = _SOURCE_ROWS,
) -> None:
    """以独立测试渲染器写入源表格的可见 PNG 投影。

    输入参数：
        path：目标 PNG 路径；rows：待可视化的 16x6 单元格值。
    输出返回值：
        无；图片包含标题、蓝色表头、12 个月份和合计行。
    """

    pytest.importorskip("PIL")
    from PIL import Image, ImageDraw, ImageFont

    widths = _COLUMN_WIDTHS
    heights = _ROW_HEIGHTS
    image = Image.new("RGB", (sum(widths), sum(heights)), "white")
    draw = ImageDraw.Draw(image)
    normal = ImageFont.load_default(size=19)
    bold = ImageFont.load_default(size=20)
    y = 0
    for row_index, (row, height) in enumerate(zip(rows, heights, strict=True)):
        x = 0
        for column_index, (value, width) in enumerate(zip(row, widths, strict=True)):
            if row_index == 2:
                fill = "#4472C4"
                foreground = "white"
            elif row_index == 15:
                fill = "#E7E6E6"
                foreground = "black"
            else:
                fill = "white"
                foreground = "black"
            if row_index not in {0, 1}:
                draw.rectangle(
                    (x, y, x + width, y + height), fill=fill, outline="black"
                )
            if value is not None:
                text = str(value)
                font = bold if row_index in {0, 2, 15} else normal
                box = draw.textbbox((0, 0), text, font=font)
                text_width = box[2] - box[0]
                text_height = box[3] - box[1]
                if row_index == 0:
                    text_x = (sum(widths) - text_width) // 2
                elif column_index == 0 or row_index == 2:
                    text_x = x + (width - text_width) // 2
                else:
                    text_x = x + width - text_width - 8
                draw.text(
                    (text_x, y + (height - text_height) // 2 - box[1]),
                    text,
                    fill=foreground,
                    font=font,
                )
            x += width
        y += height
    image.save(path, format="PNG", optimize=True)


def _write_presentation(path: Path, picture_path: Path | None) -> None:
    """写入保留正式关键文本且在第 3 页插图的合成 PPTX。

    输入参数：
        path：目标 PPTX；picture_path：已渲染的源表格图片，
            ``None`` 表示稍后由测试插入 native table。
    输出返回值：
        无；生成 5 页演示文稿，图片仅位于第 3 页内容区。
    """

    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    presentation.slide_width = 12_192_000
    presentation.slide_height = 6_858_000
    titles = (
        "Mcdonald’s annual report",
        "Annual total",
        "Here is chart of McDonald's - Regional Monthly Performance 2026",
        "Jan data",
        "March data",
    )
    bodies = (
        "2026",
        "Revenue (USD) :1398815$\nProfit (USD):316340$\nCustomers:35140",
        "",
        "Revenue (USD) :136690$\nProfit (USD):47109$\nCustomers:3602\n",
        "Revenue (USD) :81835$\nProfit (USD):-5823$\nAvg. Transaction (USD):43.16\n",
    )
    for index, (title, body) in enumerate(zip(titles, bodies, strict=True)):
        layout = (
            presentation.slide_layouts[0]
            if index == 0
            else presentation.slide_layouts[1]
        )
        slide = presentation.slides.add_slide(layout)
        title_box = slide.shapes.title
        assert title_box is not None
        title_box.text = title
        body_box = slide.placeholders[1]
        body_box.text = body
        if index == 0:
            title_box.left, title_box.top = 1_524_000, 1_122_363
            title_box.width, title_box.height = 9_144_000, 2_387_600
            body_box.left, body_box.top = 1_524_000, 3_602_038
            body_box.width, body_box.height = 9_144_000, 1_655_762
        else:
            title_box.left, title_box.top = 838_200, 365_125
            title_box.width, title_box.height = 10_515_600, 1_325_563
            body_box.left, body_box.top = 838_200, 1_825_625
            body_box.width, body_box.height = 10_515_600, 4_351_338
        if index == 2 and picture_path is not None:
            slide.shapes.add_picture(
                str(picture_path),
                2_096_000,
                2_001_294,
                width=8_000_000,
                height=4_000_000,
            )
    presentation.save(path)


def _write_native_table_presentation(path: Path) -> None:
    """写入逐格来自源 XLSX 的可编辑 PowerPoint 表格。

    输入参数：
        path：目标 PPTX 路径。
    输出返回值：
        无；第 3 页包含 16x6 表格、合并标题行及关键源样式。
    """

    _write_presentation(path, None)
    pptx = pytest.importorskip("pptx")
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    presentation = pptx.Presentation(path)
    slide = presentation.slides[2]
    table = slide.shapes.add_table(
        16,
        6,
        838_200,
        1_825_625,
        10_515_600,
        4_351_338,
    ).table
    for column, width in zip(
        table.columns,
        _allocate_fixture_dimensions(10_515_600, _COLUMN_WIDTHS),
        strict=True,
    ):
        column.width = width
    for table_row, height in zip(
        table.rows,
        _allocate_fixture_dimensions(4_351_338, _ROW_HEIGHTS),
        strict=True,
    ):
        table_row.height = height
    for row_index, row in enumerate(_SOURCE_ROWS):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = "" if value is None else str(value)
            paragraph = cell.text_frame.paragraphs[0]
            if row_index in {0, 2}:
                paragraph.alignment = PP_ALIGN.CENTER
            elif column_index == 0:
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                paragraph.alignment = PP_ALIGN.RIGHT
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = "Arial"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            if row_index in {0, 2, 15}:
                run.font.bold = True
            if row_index == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif row_index == 15:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE7, 0xE6, 0xE6)
    table.cell(0, 0).merge(table.cell(0, 5))
    presentation.save(path)


def test_source_table_picture_on_slide_three_passes(tmp_path: Path) -> None:
    """验证真实源值投影仅插入第 3 页时通过。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    _write_presentation(tmp_path / "McDonalds_powerpoint_report.pptx", picture_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is True, result
    assert result["score"] == 1.0


def test_same_rgb_png_with_extra_metadata_fails_raw_byte_identity(
    tmp_path: Path,
) -> None:
    """验证像素一致但 PNG metadata/原始字节不同时严格失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    canonical_path = tmp_path / "canonical.png"
    metadata_path = tmp_path / "same-rgb-with-metadata.png"
    _fixture_table_picture(canonical_path)
    pytest.importorskip("PIL")
    from PIL import Image, PngImagePlugin

    metadata = PngImagePlugin.PngInfo()
    metadata.add_text("fixture", "same-rgb-different-container")
    with Image.open(canonical_path) as image:
        canonical_pixels = image.convert("RGB").tobytes()
        image.save(metadata_path, format="PNG", optimize=True, pnginfo=metadata)
    with Image.open(metadata_path) as image:
        assert image.convert("RGB").tobytes() == canonical_pixels
    assert metadata_path.read_bytes() != canonical_path.read_bytes()
    _write_presentation(
        tmp_path / "McDonalds_powerpoint_report.pptx",
        metadata_path,
    )

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


@pytest.mark.parametrize(
    ("row_index", "column_index", "wrong_value"),
    [(11, 5, 48.16), (15, 1, 1_398_816)],
)
def test_picture_with_one_wrong_source_digit_fails(
    tmp_path: Path,
    row_index: int,
    column_index: int,
    wrong_value: object,
) -> None:
    """验证表格结构正确但任一数字错误时不能通过。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    rows = [list(row) for row in _SOURCE_ROWS]
    rows[row_index][column_index] = wrong_value
    picture_path = tmp_path / "wrong-source-value.png"
    _fixture_table_picture(
        picture_path,
        tuple(tuple(row) for row in rows),
    )
    _write_presentation(tmp_path / "McDonalds_powerpoint_report.pptx", picture_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


@pytest.mark.parametrize(
    ("column_index", "wrong_header"),
    [
        (1, "Pevenue (USD)"),
        (5, "Avg. Transaetion (USD)"),
    ],
)
def test_picture_with_one_wrong_source_character_fails(
    tmp_path: Path,
    column_index: int,
    wrong_header: str,
) -> None:
    """验证表头中仅一个字符漂移时确定性投影严格失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    rows = [list(row) for row in _SOURCE_ROWS]
    rows[2][column_index] = wrong_header
    picture_path = tmp_path / "wrong-source-character.png"
    _fixture_table_picture(picture_path, tuple(tuple(row) for row in rows))
    _write_presentation(
        tmp_path / "McDonalds_powerpoint_report.pptx",
        picture_path,
    )

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_native_table_matching_source_cells_passes(tmp_path: Path) -> None:
    """验证第 3 页可编辑表格逐格匹配源值时满分通过。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    _write_native_table_presentation(tmp_path / "McDonalds_powerpoint_report.pptx")

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is True, result
    assert result["score"] == 1.0


def test_native_table_grid_total_must_match_shape_extent(tmp_path: Path) -> None:
    """验证内部网格总宽与 graphicFrame 宽度不一致时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    table_shape = presentation.slides[2].shapes[-1]
    table_shape.width -= 100_000
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_native_table_compressed_source_columns_fail(tmp_path: Path) -> None:
    """验证将 B:F 列压缩为 1 EMU 且由 A 列吸收宽度时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    table = presentation.slides[2].shapes[-1].table
    total_width = sum(column.width for column in table.columns)
    for column_index in range(1, 6):
        table.columns[column_index].width = 1
    table.columns[0].width = total_width - 5
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_native_table_compressed_body_rows_fail(tmp_path: Path) -> None:
    """验证将所有月份主体行压缩为 1 EMU 且保持总高时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    table = presentation.slides[2].shapes[-1].table
    total_height = sum(row.height for row in table.rows)
    for row_index in range(3, 15):
        table.rows[row_index].height = 1
    table.rows[15].height = total_height - sum(
        table.rows[row_index].height for row_index in range(15)
    )
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_native_table_oversized_cell_margins_fail(tmp_path: Path) -> None:
    """验证单元格边距遮蔽非空文本时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    table = presentation.slides[2].shapes[-1].table
    table.cell(3, 1).margin_left = table.columns[1].width
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_correct_picture_outside_slide_three_content_area_fails(tmp_path: Path) -> None:
    """验证内容正确但移到页外的图片不能冒充可见交付。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    picture = presentation.slides[2].shapes[-1]
    picture.left = -picture.width
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_correct_picture_pixels_stretched_in_shape_fail(tmp_path: Path) -> None:
    """验证候选像素正确但形状纵横比被拉伸时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    picture = presentation.slides[2].shapes[-1]
    picture.left = 838_200
    picture.top = 1_825_625
    picture.width = 10_515_600
    picture.height = 4_351_338
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_correct_picture_blob_with_luminance_effect_fails(tmp_path: Path) -> None:
    """验证内嵌 RGB 正确但 OOXML 额外亮度效果改变显示时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    from pptx.oxml.xmlchemy import OxmlElement

    presentation = pptx.Presentation(presentation_path)
    picture = presentation.slides[2].shapes[-1]
    luminance = OxmlElement("a:lum")
    luminance.set("bright", "50000")
    picture._element.blipFill.blip.append(luminance)
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


@pytest.mark.parametrize("mutation", ["reltype", "external", "content_type"])
def test_picture_relationship_contract_mutations_fail(
    tmp_path: Path,
    mutation: str,
) -> None:
    """验证非图像类型、外链或错 ContentType 的关系均失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
    from pptx.oxml.ns import qn

    presentation = pptx.Presentation(presentation_path)
    slide = presentation.slides[2]
    picture = slide.shapes[-1]
    relationship_id = picture._element.blipFill.blip.get(qn("r:embed"))
    assert relationship_id is not None
    relationship = slide.part.rels[relationship_id]
    if mutation == "reltype":
        relationship._reltype = "urn:not-an-image"
    elif mutation == "external":
        relationship._target_mode = RTM.EXTERNAL
        relationship._target = "https://example.invalid/source-table.png"
        relationship.__dict__.pop("is_external", None)
    else:
        target_part = relationship.target_part
        target_part._content_type = "application/octet-stream"
        target_part.__dict__.pop("content_type", None)
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_picture_blip_with_spoofed_namespace_fails(tmp_path: Path) -> None:
    """验证局部名仍为 ``blip`` 但 namespace 伪造时不能通过允许列表。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    blip = presentation.slides[2].shapes[-1]._element.blipFill.blip
    blip.tag = "{urn:spoofed-drawingml}blip"
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_rgba_picture_is_rejected_even_when_rgb_projection_matches(
    tmp_path: Path,
) -> None:
    """验证候选图片携带 alpha 通道时不依赖背景合成并严格失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    rgb_path = tmp_path / "source-table-rgb.png"
    rgba_path = tmp_path / "source-table-rgba.png"
    _fixture_table_picture(rgb_path)
    pytest.importorskip("PIL")
    from PIL import Image

    with Image.open(rgb_path) as image:
        image.convert("RGBA").save(rgba_path)
    _write_presentation(
        tmp_path / "McDonalds_powerpoint_report.pptx",
        rgba_path,
    )

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_hidden_target_slide_fails(tmp_path: Path) -> None:
    """验证第 3 页被标记为隐藏时，正确插入内容仍严格失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    presentation.slides[2]._element.set("show", "0")
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_native_table_compressed_on_one_axis_fails(tmp_path: Path) -> None:
    """验证 native table 只在水平方向压缩到内容区 70% 以下时失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    table_shape = presentation.slides[2].shapes[-1]
    table_shape.left = 3_000_000
    table_shape.width = 6_000_000
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


@pytest.mark.parametrize("kind", ["blank", "random", "store1"])
def test_unrelated_picture_content_fails(tmp_path: Path, kind: str) -> None:
    """验证空白、随机或蜜雪冰城表格图片均不能通过。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    pytest.importorskip("PIL")
    from PIL import Image, ImageDraw

    picture_path = tmp_path / f"{kind}.png"
    if kind == "blank":
        Image.new("RGB", (1180, 590), "white").save(picture_path)
    elif kind == "random":
        image = Image.new("RGB", (1180, 590), "white")
        draw = ImageDraw.Draw(image)
        for y in range(0, 590, 20):
            for x in range(0, 1180, 20):
                if (x // 20 + y // 20) % 2:
                    draw.rectangle((x, y, x + 19, y + 19), fill="#D02090")
        image.save(picture_path)
    else:
        store_rows = (
            (
                "Mixue Ice Cream & Tea Chaoyang Store - 2026 Sales Statistics",
                None,
                None,
                None,
                None,
                None,
            ),
            (None, None, None, None, None, None),
            (
                "month",
                "Sales revenue(yuan)",
                "Year-on-year growth（%）",
                None,
                None,
                None,
            ),
            ("1月", 100924, -5.55, None, None, None),
            ("2月", 69917, 7.34, None, None, None),
            ("3月", 92050, 17.9, None, None, None),
            ("4月", 51583, -7.63, None, None, None),
            ("5月", 120930, 11.38, None, None, None),
            ("6月", 90228, 11.8, None, None, None),
            ("7月", 116026, 13.12, None, None, None),
            ("8月", 146386, 8.08, None, None, None),
            ("9月", 66271, 29.28, None, None, None),
            ("10月", 117566, 4.34, None, None, None),
            ("11月", 55478, -6.71, None, None, None),
            ("12月", 110618, 10.4, None, None, None),
            ("Annual total", None, None, None, None, None),
        )
        _fixture_table_picture(picture_path, store_rows)
    _write_presentation(
        tmp_path / "McDonalds_powerpoint_report.pptx",
        picture_path,
    )

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_correct_picture_on_wrong_slide_fails(tmp_path: Path) -> None:
    """验证源表格图片插入第 4 页时不能通过。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, None)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    presentation.slides[3].shapes.add_picture(
        str(picture_path),
        838_200,
        1_825_625,
        width=10_515_600,
        height=4_351_338,
    )
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_native_table_with_one_wrong_digit_fails(tmp_path: Path) -> None:
    """验证 native table 中任一单元格数字漂移时严格失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    presentation.slides[2].shapes[-1].table.cell(3, 1).text = "136691"
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_other_slide_text_tampering_fails(tmp_path: Path) -> None:
    """验证第 3 页插入正确时仍不允许破坏其他页关键内容。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    picture_path = tmp_path / "source-table.png"
    _fixture_table_picture(picture_path)
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_presentation(presentation_path, picture_path)
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation(presentation_path)
    presentation.slides[1].shapes[1].text = "tampered annual total"
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


@pytest.mark.parametrize("mutation", ["white_text", "one_point", "inherited"])
def test_native_table_unreadable_body_style_fails(
    tmp_path: Path,
    mutation: str,
) -> None:
    """验证 native table 数据格的白字或 1pt 微缩文本不能通过。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    presentation = pptx.Presentation(presentation_path)
    table = presentation.slides[2].shapes[-1].table
    for row_index in range(3, 15):
        for column_index in range(6):
            for run in (
                table.cell(row_index, column_index).text_frame.paragraphs[0].runs
            ):
                if mutation == "white_text":
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                elif mutation == "one_point":
                    run.font.size = Pt(1)
                else:
                    run.font.size = None
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


@pytest.mark.parametrize("row_index", [0, 2, 15])
def test_native_table_key_row_one_point_text_fails(
    tmp_path: Path,
    row_index: int,
) -> None:
    """验证标题、表头或合计行文本微缩到 1pt 时严格失败。"""

    _write_source_workbook(tmp_path / "McDonalds_Monthly_Data.xlsx")
    presentation_path = tmp_path / "McDonalds_powerpoint_report.pptx"
    _write_native_table_presentation(presentation_path)
    pptx = pytest.importorskip("pptx")
    from pptx.util import Pt

    presentation = pptx.Presentation(presentation_path)
    table = presentation.slides[2].shapes[-1].table
    for column_index in range(6):
        for run in table.cell(row_index, column_index).text_frame.paragraphs[0].runs:
            run.font.size = Pt(1)
    presentation.save(presentation_path)

    result = check_combinationdocs003_source_table_insert(str(tmp_path), {})

    assert result["pass"] is False, result
    assert result["score"] == 0.0


def test_pinned_contract_exactly_matches_formal_input_manifest() -> None:
    """验证 003 evaluator 的路径、大小、摘要和改动策略与 manifest 一致。"""

    task = _load_task()
    contract = operation_evaluator_module._PINNED_ARTIFACT_CONTRACTS[_TASK_ID]
    manifest_path = _REPO_ROOT / str(task["asset_manifest"])
    payload = manifest_path.read_bytes()
    manifest = json.loads(payload)

    assert contract.manifest_reference == task["asset_manifest"]
    assert hashlib.sha256(payload).hexdigest() == contract.manifest_sha256
    assert tuple(
        (entry["path"], entry["size"], entry["sha256"]) for entry in manifest["files"]
    ) == tuple((item.path, item.size, item.sha256) for item in contract.files)
    assert tuple(item.preserved for item in contract.files) == (
        True,
        False,
        True,
        True,
    )
    assert tuple(item.must_change for item in contract.files) == (
        False,
        True,
        False,
        False,
    )
    assert contract.expected_document_count == 1


@pytest.mark.parametrize("ppt_identity", ["no_op", "unreadable"])
def test_pinned_contract_rejects_noop_or_unreadable_output_identity(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    ppt_identity: str,
) -> None:
    """验证 PPT 与原输入字节相同或身份读取失败时 fail closed。"""

    task = _load_task()
    contract = operation_evaluator_module._PINNED_ARTIFACT_CONTRACTS[_TASK_ID]
    paths = tuple(tmp_path / item.path for item in contract.files)
    for path in paths:
        path.write_bytes(b"synthetic-contract-path")
    expected_by_name = {item.path: (item.size, item.sha256) for item in contract.files}

    def fake_identity(path: Path) -> tuple[int, str] | None:
        """为合同分支返回受控的文件身份。"""

        if (
            path.name == "McDonalds_powerpoint_report.pptx"
            and ppt_identity == "unreadable"
        ):
            return None
        return expected_by_name[path.name]

    monkeypatch.setattr(
        operation_evaluator_module,
        "_regular_file_identity",
        fake_identity,
    )

    result = operation_evaluator_module._pinned_artifact_contract_failure(
        OPERATION_TASK_RULES[_TASK_ID],
        task,
        tmp_path,
        paths,
    )

    assert result is not None
    assert result.passed is False
    assert result.score == 0.0
    assert result.reason_codes == ("ARTIFACT_CONTRACT_MISMATCH",)
    assert result.rule_results[0].evaluated_artifact_count == 1


def test_pinned_contract_accepts_only_changed_ppt_with_three_preserved_xlsx(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """验证三份 XLSX 身份不变且 PPT 已改动时合同门才放行。"""

    task = _load_task()
    contract = operation_evaluator_module._PINNED_ARTIFACT_CONTRACTS[_TASK_ID]
    paths = tuple(tmp_path / item.path for item in contract.files)
    for path in paths:
        path.write_bytes(b"synthetic-contract-path")
    expected_by_name = {item.path: (item.size, item.sha256) for item in contract.files}

    def fake_identity(path: Path) -> tuple[int, str]:
        """返回三份 pinned 源文件与一份已改 PPT 的身份。"""

        if path.name == "McDonalds_powerpoint_report.pptx":
            return (41_100, "0" * 64)
        return expected_by_name[path.name]

    monkeypatch.setattr(
        operation_evaluator_module,
        "_regular_file_identity",
        fake_identity,
    )

    result = operation_evaluator_module._pinned_artifact_contract_failure(
        OPERATION_TASK_RULES[_TASK_ID],
        task,
        tmp_path,
        paths,
    )

    assert result is None


@pytest.mark.parametrize(
    "tampered_name",
    ["McDonalds_Monthly_Data.xlsx", "store1.xlsx", "store2.xlsx"],
)
def test_pinned_contract_rejects_each_tampered_xlsx(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    tampered_name: str,
) -> None:
    """验证三份 XLSX 中任一字节身份漂移都失败。"""

    task = _load_task()
    contract = operation_evaluator_module._PINNED_ARTIFACT_CONTRACTS[_TASK_ID]
    paths = tuple(tmp_path / item.path for item in contract.files)
    for path in paths:
        path.write_bytes(b"synthetic-contract-path")
    expected_by_name = {item.path: (item.size, item.sha256) for item in contract.files}

    def fake_identity(path: Path) -> tuple[int, str]:
        """仅对参数指定的 XLSX 返回漂移身份。"""

        if path.name == tampered_name:
            return (1, "f" * 64)
        if path.name == "McDonalds_powerpoint_report.pptx":
            return (41_100, "0" * 64)
        return expected_by_name[path.name]

    monkeypatch.setattr(
        operation_evaluator_module,
        "_regular_file_identity",
        fake_identity,
    )

    result = operation_evaluator_module._pinned_artifact_contract_failure(
        OPERATION_TASK_RULES[_TASK_ID],
        task,
        tmp_path,
        paths,
    )

    assert result is not None
    assert result.reason_codes == ("ARTIFACT_CONTRACT_MISMATCH",)


@pytest.mark.parametrize("mutation", ["missing", "extra", "renamed"])
def test_pinned_contract_rejects_nonexact_four_file_path_closure(
    tmp_path: Path,
    mutation: str,
) -> None:
    """验证四文件闭集中任一缺失、额外或改名都固定失败。"""

    task = _load_task()
    contract = operation_evaluator_module._PINNED_ARTIFACT_CONTRACTS[_TASK_ID]
    paths = [tmp_path / item.path for item in contract.files]
    if mutation == "missing":
        paths.pop()
    elif mutation == "extra":
        paths.append(tmp_path / "extra.xlsx")
    else:
        paths[-1] = tmp_path / "renamed.xlsx"
    for path in paths:
        path.write_bytes(b"synthetic-contract-path")

    result = operation_evaluator_module._pinned_artifact_contract_failure(
        OPERATION_TASK_RULES[_TASK_ID],
        task,
        tmp_path,
        tuple(sorted(paths, key=lambda path: path.name)),
    )

    assert result is not None
    assert result.reason_codes == ("ARTIFACT_CONTRACT_MISMATCH",)
    assert result.rule_results[0].evaluated_artifact_count == 1
